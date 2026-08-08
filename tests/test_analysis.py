import os
import tempfile
import unittest

import pandas as pd

from models.analysis import (
    analyze_status,
    analyze_status_derived,
    analyze_bottleneck_machines,
    analyze_steps,
    analyze_steps_sa,
    build_cycles_df,
    build_gantt_rows,
    build_gantt_rows_sa,
    detect_units_per_tray,
    detect_tray_stats,
    down_pareto,
    measure_tray_change,
    parse_em_production,
    summarize_alarms,
    summarize_eff_coretech,
    summarize_uph,
    summarize_uph_ame,
)
from models.log_model import LogModel


def row(content, fname="a.log"):
    return {"FileName": fname, "Content": content}


CYCLE_ROWS = [
    row("2026-07-08 02:30:29.020   动作时间:02:30:29.006;模块名称:取料2(工位7);动作名称:放生料完成;;;;"),
    row("2026-07-08 02:30:39.020   动作时间:02:30:39.006;模块名称:取料2(工位7);动作名称:放生料完成;;;;"),
    row("2026-07-08 02:30:59.020   动作时间:02:30:59.006;模块名称:取料2(工位7);动作名称:放生料完成;;;;"),
    row("2026-07-08 02:35:00.000   动作时间:02:35:00;模块名称:取料2(工位7);动作名称:放生料完成;;;;"),
]


class AnalysisTest(unittest.TestCase):
    def test_detect_tray_stats_line_mode(self):
        # 上料机：每条 carrier 消息 = 一盘，每盘 8 颗（SiteId 数量）
        rows = []
        t = 0.0
        for i in range(5):
            sites = "".join('"SiteId":%d,' % j for j in range(1, 9))
            rows.append(row(
                "%02d:%02d:%02d.%03d   Info <MES> 更新Carrier盘:carrierId:TACA%08d,Data:[%s]" % (
                    int(t // 3600), int(t % 3600 // 60), int(t % 60), int(t % 1 * 1000), i, sites),
                "上料機/Logger.txt",
            ))
            t += 26.0
        stats = detect_tray_stats(rows, r"更新Carrier盘:.*?carrierId:([A-Z0-9\-]+)",
                                  r"SiteId", segments="line")
        self.assertIsNotNone(stats)
        self.assertEqual(stats["units_per_tray"], 8)
        self.assertEqual(stats["tray_count"], 4)
        self.assertAlmostEqual(stats["tray_seconds"], 26.0, delta=0.1)

    def test_detect_tray_stats_id_mode_and_run_split(self):
        # 下料机：CubeTrayId 每盘 24 颗；同一 ID 跨小时复用要按 run_gap 切段
        rows = []
        for hour in range(2):
            base = hour * 3600
            for tray in ("TCP-A", "TCP-B"):
                for k in range(24):
                    rows.append(row(
                        "%02d:%02d:%02d.%03d   Info <StepPrTrayCell> 获取格子上料机;===>CubeTrayId:%s,CubeTrayCell:%d" % (
                            base // 3600, (base + k * 3) % 3600 // 60, (base + k * 3) % 60,
                            (base + k * 3) % 1 * 1000, tray, k + 1),
                        "下料機/Logger.txt",
                    ))
        stats = detect_tray_stats(rows, r"CubeTrayId:([A-Z0-9\-]+)", segments="id", run_gap=300.0)
        self.assertIsNotNone(stats)
        self.assertEqual(stats["units_per_tray"], 24)
        # TCP-A 在两个小时内各出现一次 -> 2 段；共 4 段，换盘 3 次
        self.assertEqual(stats["segments"], 4)
        self.assertEqual(stats["tray_count"], 3)

    def test_measure_tray_change(self):
        # SA 式换盘：卸载 -> 下一次装载 的间隔中位数
        rows = []
        for i in range(4):
            unload = 10.0 + i * 100.0
            load = unload + 12.0
            rows.append(row("%02d:%02d:%02d.%03d  Info 清除2号托盘所有格子状态" % (
                int(unload // 3600), int(unload % 3600 // 60), int(unload % 60),
                int(unload % 1 * 1000)), "下料機/Logger.txt"))
            rows.append(row("%02d:%02d:%02d.%03d  Error 轨道2进板成功" % (
                int(load // 3600), int(load % 3600 // 60), int(load % 60),
                int(load % 1 * 1000)), "下料機/Logger.txt"))
        stats = measure_tray_change(rows, "清除2号托盘所有格子状态", "轨道2进板成功")
        self.assertIsNotNone(stats)
        self.assertEqual(stats["tray_count"], 4)
        self.assertAlmostEqual(stats["tray_seconds"], 12.0, delta=0.1)

    def test_detect_units_per_tray(self):
        # LM 式：CCD 每批定位多颗（MarkEnd1），每盘多批；批次/批数由日志自动统计
        def mk(ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "LM/l.txt",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})

        rows = []
        for tray, base in enumerate([0.0, 40.0]):
            mk(base, "Move to unload")
            for b in range(3):
                ccd = base + 1.0 + b * 10.0
                mk(ccd, "Start parsing CCD data")
                for k in range(4):
                    mk(ccd + 0.5 + k * 1.0, "MarkEnd1")
        stats = detect_units_per_tray(rows, "Start parsing CCD data", "MarkEnd1", "Move to unload")
        self.assertEqual(stats["units_per_batch"], 4)   # 每批 4 颗
        self.assertEqual(stats["batches_per_tray"], 3)  # 每盘 3 批
        self.assertEqual(stats["units_per_tray"], 12)   # 每盘 12 颗

    def test_summarize_uph_tray_overhead(self):
        # 换盘时间平摊到整盘产品：有效周期 = 正常周期平均 + 每颗换盘开销
        df = build_cycles_df(CYCLE_ROWS, "放生料完成,放熟料完成")
        summary = summarize_uph(df, 1, tray_overhead_seconds=0.5)
        row0 = summary.iloc[0]
        self.assertEqual(row0['每颗换盘开销(秒)'], 0.5)
        self.assertEqual(row0['有效周期(秒)'], 10.5)
        self.assertEqual(row0['Pure UPH(个/小时)'], round(3600.0 / 10.5, 2))
        ame = summarize_uph_ame(df, 1, tray_overhead_seconds=0.5)
        self.assertEqual(ame.iloc[0]['Pure UPH(个/小时)'], round(3600.0 / 10.5, 2))

    def test_analyze_steps_anomaly_and_ignore_fast(self):
        def mk(rows, ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "LM/l.txt",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})

        rows = []
        # 3 个循环：GetSN(0.2/0.18/1.2s)、CCD定位(0.002s 快步骤)、打标(0.3/0.28/2.5s)
        for i, (gs, mark) in enumerate([(0.20, 0.30), (0.18, 0.28), (1.20, 2.50)]):
            base = i * 10.0
            mk(rows, base, "GetSN_Start")
            mk(rows, base + gs, "GetSN_End")
            mk(rows, base + gs + 0.002, "Ccd Locate Data Check success")
            mk(rows, base + gs + 0.3, "MarkEnd1")

        units = [{
            "name": "整机", "cycle": "MarkEnd1",
            "steps": [
                {"name": "GetSN", "start": "GetSN_Start", "end": "GetSN_End",
                 "timeout_seconds": 0.5},
                {"name": "CCD定位", "end": "Ccd Locate Data Check success"},
                {"name": "打标", "end": "MarkEnd1"},
            ],
        }]
        df = analyze_steps(rows, units, coefficient=1.5)
        steps = set(df["步骤"])
        self.assertNotIn("CCD定位", steps)          # 中位 0.002s < 0.01s 被忽略
        self.assertIn("GetSN", steps)
        gsn = df[df["步骤"] == "GetSN"].iloc[0]
        self.assertEqual(gsn["异常次数"], 1)         # 1.2s > 0.69+0.5
        self.assertEqual(gsn["中位时长"], "0.69")
        self.assertEqual(gsn["异常影响时长"], "0.51")
        mark = df[df["步骤"] == "打标"].iloc[0]
        self.assertGreaterEqual(mark["循环数"], 1)

    def test_fmt_duration_rule(self):
        # 约定：<60 秒按秒显示，>=60 秒按 h:mm:ss 显示
        from models.analysis import fmt_duration
        self.assertEqual(fmt_duration(0.05), "0.05")
        self.assertEqual(fmt_duration(0.1), "0.1")
        self.assertEqual(fmt_duration(1.3), "1.3")
        self.assertEqual(fmt_duration(7.5), "7.5")
        self.assertEqual(fmt_duration(45), "45")
        self.assertEqual(fmt_duration(59.999), "59.999")
        self.assertEqual(fmt_duration(60), "0:01:00")
        self.assertEqual(fmt_duration(61.5), "0:01:01.500")
        self.assertEqual(fmt_duration(86396.148), "23:59:56.148")
        # PPT 汇总展示：h:mm:ss 只到秒
        self.assertEqual(fmt_duration(61.5, with_ms=False), "0:01:01")
        self.assertEqual(fmt_duration(86396.148, with_ms=False), "23:59:56")
        self.assertEqual(fmt_duration(""), "")

    def test_build_gantt_rows(self):
        def mk(rows, ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "LM/l.txt",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})

        rows = []
        # 3 个周期，每周期 4s：A 步骤 [0,1]，B 步骤 [1,3]
        for i in range(3):
            base = i * 5.0
            mk(rows, base, "MarkEnd1")
            mk(rows, base + 1.0, "StepA_Done")
            mk(rows, base + 3.0, "StepB_Done")
        units = [{
            "name": "整机", "cycle": "MarkEnd1",
            "steps": [
                {"name": "A", "end": "StepA_Done"},
                {"name": "B", "end": "StepB_Done"},
            ],
        }]
        df = build_gantt_rows(rows, units)
        self.assertEqual(list(df["步骤"]), ["A", "B"])   # 按开始秒排序
        a = df[df["步骤"] == "A"].iloc[0]
        b = df[df["步骤"] == "B"].iloc[0]
        self.assertEqual(a["开始秒"], 0.0)
        self.assertEqual(a["结束秒"], 1.0)
        self.assertEqual(b["开始秒"], 1.0)
        self.assertEqual(b["结束秒"], 3.0)
        self.assertEqual(list(df["层级"]), ["循环", "循环"])

    def test_analyze_steps_exclude_long_stop(self):
        def mk(rows, ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "FR/l.txt",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})

        rows = []
        # 首边界后 3 个循环：正常 3s、超长停机 2000s（>900 应剔除）、异常慢 12s
        mk(rows, 0.0, "左轴点胶完成")
        base = 1.0
        for i, dur in enumerate([3.0, 2000.0, 12.0]):
            mk(rows, base, "左轴开始点胶")
            base += dur
            mk(rows, base, "左轴点胶完成")
            base += 1.0

        units = [{
            "name": "左轴", "module": {"pattern": "左轴"}, "cycle": "轴点胶完成",
            "steps": [{"name": "点胶", "start": "左轴开始点胶", "end": "左轴点胶完成"}],
        }]
        df = analyze_steps(rows, units, coefficient=1.5, max_step_seconds=900.0)
        row0 = df.iloc[0]
        self.assertEqual(row0["循环数"], 2)          # 2000s 的停机循环被剔除
        self.assertEqual(row0["异常次数"], 1)          # 只剩 12s 这一个异常
        self.assertEqual(row0["中位时长"], "7.5")
        self.assertEqual(row0["异常影响时长"], "4.5")  # 12 − 7.5

    def test_analyze_steps_standalone_keeps_chain_partition(self):
        """standalone 步骤独立计时但不打断链式切分：单颗循环段和仍等于周期长。"""
        def mk(rows, ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "LM/l.txt",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})

        rows = []
        # 周期 = MarkEnd1 间隔 1.4s：打标 1.3s + 间隔 0.1s；读码每 2 个周期插一次（standalone）
        for i in range(6):
            base = i * 1.4
            if i % 2 == 0:
                mk(rows, base + 0.05, "GetSN_Start")
                mk(rows, base + 0.35, "GetSN_End")
            mk(rows, base + 0.4, "MarkStart1")
            mk(rows, base + 1.7, "MarkEnd1")

        units = [{
            "name": "整机", "cycle": "MarkEnd1",
            "steps": [
                {"name": "打标前间隔", "end": "MarkStart1"},
                {"name": "打标", "start": "MarkStart1", "end": "MarkEnd1"},
                {"name": "读码(standalone)", "start": "GetSN_Start", "end": "GetSN_End",
                 "standalone": True},
            ],
        }]
        df = analyze_steps(rows, units, coefficient=1.5)
        gap = df[df["步骤"] == "打标前间隔"].iloc[0]
        mark = df[df["步骤"] == "打标"].iloc[0]
        read = df[df["步骤"] == "读码(standalone)"].iloc[0]
        # standalone 不参与链式切分：打标前间隔始终=0.1s（不被 GetSN_End 锚点打断）
        self.assertEqual(gap["中位时长"], "0.1")
        self.assertEqual(mark["中位时长"], "1.3")
        # 单颗循环段和 = 周期长 1.4s
        def to_sec(s):
            if ":" not in s:
                return float(s)
            hh, mm, ss = s.split(":")
            return int(hh) * 3600 + int(mm) * 60 + float(ss)
        self.assertAlmostEqual(to_sec(gap["中位时长"]) + to_sec(mark["中位时长"]), 1.4, delta=0.01)
        self.assertEqual(read["中位时长"], "0.3")

    def test_analyze_steps_sa_stations(self):
        def mk(rows, ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "SA/l.txt",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})

        rows = []
        t = 0.0
        # 4 排：每排 点胶→贴附×2→热压→检测×2；热压间隔 10/10/10/30s（末排异常）
        for r, hp in enumerate([10.0, 10.0, 10.0, 30.0]):
            mk(rows, t, "DispOneChipProfileWorkCycle")
            mk(rows, t + 2.0, "AfterPickUp StopCondition")
            mk(rows, t + 3.0, "AfterPickUp StopCondition")
            mk(rows, t + hp, "Heater 0 :Heating Complete")
            mk(rows, t + hp + 0.5, "UDP Module - Good")
            mk(rows, t + hp + 1.0, "UDP Module - Good")
            t += hp

        df = analyze_steps_sa(rows, coefficient=1.5)
        steps = set(df["步骤"])
        self.assertEqual(steps, {"点胶", "贴附", "热压", "检测"})
        hp = df[df["步骤"] == "热压"].iloc[0]
        self.assertEqual(hp["循环数"], 3)
        self.assertEqual(hp["异常次数"], 1)          # 30s > 10×1.5
        self.assertEqual(hp["中位时长"], "10")
        self.assertEqual(hp["异常影响时长"], "20")
        dp = df[df["步骤"] == "点胶"].iloc[0]
        self.assertEqual(dp["循环数"], 3)
        self.assertEqual(dp["异常次数"], 0)
        jc = df[df["步骤"] == "检测"].iloc[0]
        self.assertEqual(jc["循环数"], 3)

    def test_analyze_steps_sa_dispense_phases(self):
        # 右点胶头内部动作：视觉对位→探针对位→点胶轮廓，B 模式段和≈该头行周期
        def mk(ts, tag, head="SeqCycle005_RightDispenserPart.cs"):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "SA/l.txt",
                         "Content": "[000] Sequence, , 20260803-%02d-%02d-%02d-%03d, %s, 30, %s"
                                    % (h, m, sec, ms, head, tag)})

        rows = []
        for base in (0.0, 20.0, 40.0):
            mk(base, "DispOneChipProfileWorkCycle()")
            mk(base + 5.0, "DispOneChipAlignVisionCycle()")
            mk(base + 12.0, "DispOneChipProbeAlignCycle()")
        df = analyze_steps_sa(rows, coefficient=1.5)
        for phase, expect in (("视觉对位", "5"), ("探针对位", "7"),
                              ("点胶轮廓", "8")):
            r = df[(df["单元"] == "右头") & (df["步骤"] == phase)].iloc[0]
            self.assertEqual(r["中位时长"], expect)
            self.assertEqual(r["循环数"], 2)

    def test_build_gantt_rows_sa(self):
        def mk(ts, tag, head="SeqCycle005_RightDispenserPart.cs"):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "SA/l.txt",
                         "Content": "[000] Sequence, , 20260803-%02d-%02d-%02d-%03d, %s, 30, %s"
                                    % (h, m, sec, ms, head, tag)})

        rows = []
        for base in (0.0, 30.0, 60.0):
            mk(base, "DispOneChipProfileWorkCycle()")
            mk(base + 5.0, "DispOneChipAlignVisionCycle()")
            mk(base + 12.0, "DispOneChipProbeAlignCycle()")
            mk(base + 20.0, "Heater 0 :Heating Complete", "SeqCycle009_HeatIndex.cs")
            mk(base + 20.5, "AfterPickUp StopCondition", "SeqCycle007_AttachPart.cs")
            mk(base + 21.0, "AfterPickUp StopCondition", "SeqCycle007_AttachPart.cs")
            mk(base + 22.0, "UDP Module - Good", "SeqCycle014_InspectionIndex.cs")
        df = build_gantt_rows_sa(rows)
        hp = df[(df["单元"] == "热压") & (df["步骤"] == "行周期")].iloc[0]
        self.assertEqual(hp["结束秒"], 30.0)            # 热压行周期 0→30
        vis = df[(df["单元"] == "右头") & (df["步骤"] == "视觉对位")].iloc[0]
        probe = df[(df["单元"] == "右头") & (df["步骤"] == "探针对位")].iloc[0]
        prof = df[(df["单元"] == "右头") & (df["步骤"] == "点胶轮廓")].iloc[0]
        self.assertEqual(vis["开始秒"], 0.0)
        self.assertEqual(vis["结束秒"], 5.0)
        self.assertEqual(probe["开始秒"], 5.0)
        self.assertEqual(probe["结束秒"], 12.0)
        self.assertEqual(prof["开始秒"], 12.0)
        self.assertEqual(prof["结束秒"], 30.0)

    def test_gantt_page_rows_picks_median_slide(self):
        # CAW 瓶颈=焊接机、4 个滑台：按周期排序取表现居中的滑台（周期 10/20/30/40 → 滑台3）
        import os
        import tempfile
        from models.report import _gantt_page_rows
        fd, path = tempfile.mkstemp(suffix=".xlsx")
        os.close(fd)
        try:
            rows = []
            for n, cyc in ((1, 10.0), (2, 20.0), (3, 30.0), (4, 40.0)):
                for side in ("左", "右"):
                    rows.append({"单元": "焊接机-滑台%d%s" % (n, side),
                                 "步骤": "下熟料", "开始秒": 0.0,
                                 "结束秒": cyc, "时长秒": cyc, "层级": "循环"})
            with pd.ExcelWriter(path, engine="xlsxwriter") as w:
                pd.DataFrame(rows).to_excel(w, sheet_name="步骤甘特图", index=False)
                pd.DataFrame({"瓶颈机台": ["焊接机"], "瓶颈CT(秒)": [2.88],
                              "UPH(个/小时)": [1250]}).to_excel(
                    w, sheet_name="AMESummary", index=False)
            sel, title = _gantt_page_rows(path)
            self.assertIn("滑台3", title)
            self.assertEqual(set(sel["单元"]), {"焊接机-滑台3左", "焊接机-滑台3右"})
        finally:
            if os.path.exists(path):
                os.remove(path)

    def test_analyze_bottleneck_machines(self):
        def mk(rows, ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "CAW/记录PLC2当前工位.log",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})

        rows = []
        # 上料机取料1/取料2：每 10s 完成一循环（4颗/循环）
        for t in range(0, 100, 10):
            mk(rows, t, "取料1 放熟料完成")
            mk(rows, t, "取料2 放熟料完成")
        # 焊接机滑台1左/右：每 20s 一循环（2颗/循环）
        for t in range(0, 100, 20):
            mk(rows, t, "滑台1左 去交换位")
            mk(rows, t, "滑台1右 去交换位")

        machines = [
            {"name": "上料机", "module": {"file": "记录PLC2"},
             "trigger": "放熟料完成", "units_per_cycle": 4,
             "unit_pattern": "取料[12]"},
            {"name": "焊接机", "module": {"file": "记录PLC2"},
             "trigger": "去交换位", "units_per_cycle": 2,
             "unit_pattern": "滑台([1-9](?:左|右))"},
        ]
        df, bn, cycles = analyze_bottleneck_machines(rows, machines)
        self.assertEqual(df.set_index("机台").loc["上料机", "单颗CT(秒)"], 2.5)    # 10/4
        self.assertEqual(df.set_index("机台").loc["焊接机", "单颗CT(秒)"], 10.0)   # 20/2
        self.assertEqual(bn["瓶颈机台"], "焊接机")
        self.assertEqual(bn["UPH(个/小时)"], 360.0)  # 3600/10
        self.assertFalse(cycles.empty)
        self.assertIn("Class", cycles.columns)

    def test_highlight_abnormal_rows(self):
        import os as _os
        from openpyxl import load_workbook
        df = pd.DataFrame({
            "FileName": ["a.log", "a.log", "a.log"],
            "Content": ["正常日志行", "解析异常 Newtonsoft.Json", "RDA620700NG55423C 产品码"],
        })
        tmp = tempfile.mkdtemp(prefix="hl_")
        path = _os.path.join(tmp, "t.xlsx")
        LogModel._write_sheets(path, {"AllLogs": df})
        # 调用方已剔除 NG（产品码含 NG 会误标），此处模拟剔除后的关键词集
        LogModel._highlight_abnormal_rows(path, "报警,异常,失败,AutoRun Stop,超时")
        wb = load_workbook(path)
        ws = wb["AllLogs"]
        reds = []
        for row in ws.iter_rows(min_row=2):
            red = bool(row[1].fill and row[1].fill.fgColor and row[1].fill.fgColor.rgb == '00FFC7CE')
            reds.append(red)
        # 异常行标红、正常行与 NG 产品码行不标红（NG 由调用方剔除；这里传入 NG 会命中，故用第二行验证）
        self.assertTrue(reds[1])
        self.assertFalse(reds[0])
        self.assertFalse(reds[2])  # NG 产品码行：调用方剔除 NG 后不标红

    def test_highlight_step_lines_and_anomaly_lines(self):
        import os as _os
        from openpyxl import load_workbook
        # analyze_steps 应跨单元收集异常触发行（不只最后一个单元）
        def mk(rows, ts, tag):
            h = int(ts // 3600); m = int(ts % 3600 // 60)
            s = ts % 60
            sec = int(s); ms = int(round((s - sec) * 1000))
            rows.append({"FileName": "CAW/记录PLC2.log",
                         "Content": "%02d:%02d:%02d.%03d %s" % (h, m, sec, ms, tag)})
        rows = []
        t = 0.0
        for r in range(4):
            mk(rows, t, "取料1 放熟料完成")
            mk(rows, t + 2.0, "取料2 放熟料完成")
            t += 40.0 if r == 2 else 10.0   # 第 3→4 循环间隔 40s → 异常
        units = [
            {"name": "取料1", "module": {"pattern": "取料1"}, "cycle": "放熟料完成",
             "steps": [{"name": "放料", "end": "放熟料完成"}]},
            {"name": "取料2", "module": {"pattern": "取料2"}, "cycle": "放熟料完成",
             "steps": [{"name": "放料", "end": "放熟料完成"}]},
        ]
        df = analyze_steps(rows, units, coefficient=1.5)
        lines = df.attrs.get("anomaly_lines") or set()
        self.assertEqual(len(lines), 2)   # 取料1/取料2 各一条异常触发行
        # 标红：写入 Excel 后按行内容精确匹配
        tmp = tempfile.mkdtemp(prefix="sl_")
        path = _os.path.join(tmp, "t.xlsx")
        LogModel._write_sheets(path, {"AllLogs": pd.DataFrame({"FileName": [r["FileName"] for r in rows],
                                                                "Content": [r["Content"] for r in rows]})})
        LogModel._highlight_step_lines(path, lines)
        wb = load_workbook(path)
        ws = wb["AllLogs"]
        reds = sum(1 for row in ws.iter_rows(min_row=2)
                   if row[1].fill and row[1].fill.fgColor and row[1].fill.fgColor.rgb == '00FFC7CE')
        self.assertEqual(reds, len(lines))

    def test_iter_monotonic_time_only_wrap(self):
        from models.analysis import _iter_monotonic
        # 纯时间日志跨零点：00:00 之后自动累加一天，保证时间单调
        rows = [
            row("23:59:50.000  A"),
            row("23:59:55.000  B"),
            row("00:00:05.000  C"),
            row("00:00:10.000  D"),
            row("无时间戳的行"),
        ]
        out = list(_iter_monotonic(rows))
        self.assertEqual(len(out), 4)
        self.assertAlmostEqual(out[0][0], 86390.0)
        self.assertAlmostEqual(out[1][0], 86395.0)
        self.assertAlmostEqual(out[2][0], 86405.0)  # 5 + 86400
        self.assertAlmostEqual(out[3][0], 86410.0)
        self.assertEqual(out[2][1], "00:00:05.000  C")

    def test_iter_monotonic_dated_rows(self):
        from models.analysis import _iter_monotonic
        rows = [
            row("2026-08-05 00:00:00.000  A"),
            row("2026-08-05 00:00:10.000  B"),
            row("2026-08-05 00:00:25.000  C"),
        ]
        out = list(_iter_monotonic(rows))
        self.assertEqual(len(out), 3)
        self.assertAlmostEqual(out[1][0] - out[0][0], 10.0)
        self.assertAlmostEqual(out[2][0] - out[1][0], 15.0)

    def test_row_date_from_filename(self):
        from models.analysis import _row_date
        self.assertIsNotNone(_row_date("上料機/Logger 2026_08_05_00.txt"))
        self.assertIsNotNone(_row_date("主機/OHLog20260805000147.txt"))
        self.assertIsNone(_row_date("LM log/plain.txt"))

    def test_status_derived_acf_stops(self):
        # ACF 主機 Err=NNNN / 下料機 生产流程出现异常 都应归类为 DOWN
        from models.analysis import analyze_status_derived
        rows = [
            row("00:00:00.000  Info <StepPrTrayCell> UnloadDuts Finish", "下料機/Logger 2026_08_05_00.txt"),
            row("00:00:10.000  Info 生产流程出现异常：正在中止线程", "下料機/Logger 2026_08_05_00.txt"),
            row("2026-08-05 00:00:20.000  ErrJob,ErrOn,Err=3539:Process Data Sending Error", "主機/OHLog20260805000001.txt"),
            row("00:00:30.000  Info UnloadDuts Finish", "下料機/Logger 2026_08_05_00.txt"),
        ]
        summary, _, detail = analyze_status_derived(
            rows, "UnloadDuts Finish", "ErrOn,生产流程出现异常,当前设备状态Maunal")
        states = set(detail["Status"])
        self.assertIn("DOWN", states)
        self.assertEqual(summary.set_index("状态").loc["DOWN", "总时长(秒)"] > 0, True)

    def test_parse_ts(self):
        from models.analysis import parse_ts
        ts = parse_ts("2026-07-08 02:30:29.020   动作时间...")
        self.assertIsNotNone(ts)
        self.assertEqual(ts.hour, 2)
        self.assertEqual(ts.second, 29)
        bracket = parse_ts("[00:00:05] 左 图像数据获取完成")
        self.assertIsNotNone(bracket)
        self.assertEqual((bracket.hour, bracket.minute, bracket.second), (0, 0, 5))
        sa = parse_ts("[001] Sequence, VisionAlign, 20260803-09-55-26-191, DispProbeCheck, 13, Head No = 0")
        self.assertIsNotNone(sa)
        self.assertEqual((sa.year, sa.month, sa.day, sa.hour, sa.minute, sa.second, sa.microsecond),
                         (2026, 8, 3, 9, 55, 26, 191000))
        acf = parse_ts("2026/07/17 18:35:49.109,127.0.0.1,H,cms2100_trace,ProcessJob,Cavity cnt:1")
        self.assertIsNotNone(acf)
        self.assertEqual((acf.year, acf.month, acf.day, acf.hour, acf.minute, acf.second),
                         (2026, 7, 17, 18, 35, 49))
        acf_ms = parse_ts("00:00:02 079    [89]    Info    <PlaceDut2Socket>    吸嘴【1】开始夹具【1】放料")
        self.assertIsNotNone(acf_ms)
        self.assertEqual(acf_ms.microsecond, 79000)

    def test_trigger_boundary(self):
        # MarkEnd1 不应匹配 MarkEnd1_0
        rows = [
            row("2026-07-08 00:00:01.000 MarkEnd1"),
            row("2026-07-08 00:00:02.000 MarkEnd1_0"),
            row("2026-07-08 00:00:04.000 MarkEnd1"),
        ]
        df = build_cycles_df(rows, "MarkEnd1")
        self.assertEqual(len(df), 1)
        self.assertEqual(df.iloc[0]['CycleSeconds'], 3.0)

    def test_phrase_trigger_with_spaces(self):
        # 含空格的完成动作短语应整体匹配（如早期 SA 的 UDP Module - Good）
        rows = [
            {"FileName": "a.txt", "Content": "[324] Sequence, VisionAlign, 20260803-09-59-09-943, X, 0, UDP Module - Good, 0"},
            {"FileName": "a.txt", "Content": "[325] Sequence, VisionAlign, 20260803-09-59-11-394, X, 0, UDP Module - Good, 0"},
            {"FileName": "a.txt", "Content": "[326] Sequence, X, 20260803-09-59-13-000, X, 0, UDP Recv [Good]"},
        ]
        df = build_cycles_df(rows, "UDP Module - Good")
        self.assertEqual(len(df), 1)  # 只匹配完整短语
        self.assertEqual(df.iloc[0]['CycleSeconds'], 1.451)

    def test_cycle_classification(self):
        df = build_cycles_df(CYCLE_ROWS, "放生料完成,放熟料完成")
        self.assertEqual(len(df), 3)
        self.assertEqual(list(df['Class']), ["正常周期", "异常周期", "异常周期"])

    def test_midnight_rollover(self):
        rows = [
            row("23:59:58.000   动作名称:放生料完成;;;"),
            row("00:00:02.000   动作名称:放生料完成;;;"),
        ]
        df = build_cycles_df(rows, "放生料完成")
        self.assertEqual(df.iloc[0]['CycleSeconds'], 4.0)

    def test_uph_summary(self):
        df = build_cycles_df(CYCLE_ROWS, "放生料完成,放熟料完成")
        summary = summarize_uph(df, units_per_cycle=1)
        row0 = summary.iloc[0]
        self.assertEqual(row0['周期总数'], 3)
        expected = round(3 / (df['CycleSeconds'].sum() / 3600.0), 2)
        self.assertEqual(row0['UPH(个/小时)'], expected)
        # Pure UPH 基于正常周期：3600/10 = 360
        self.assertEqual(row0['Pure UPH(个/小时)'], 360.0)
        # 无理想/最大CT时，Derated M2 = 全部周期平均
        self.assertEqual(row0['Derated UPH M2(个/小时)'], expected)

    def test_uph_derated_m2_filters(self):
        df = build_cycles_df(CYCLE_ROWS, "放生料完成,放熟料完成")
        summary = summarize_uph(df, 1, ideal_ct=10.0, max_ct=20.0)
        row0 = summary.iloc[0]
        self.assertEqual(row0['Pure UPH(个/小时)'], 360.0)        # 3600 / 10
        self.assertEqual(row0['Derated UPH M2(个/小时)'], 240.0)   # 剔除 241 后 avg=15 → 3600/15

    def test_pure_uph_factor(self):
        # FR 并行工位：单轴 Pure UPH ×0.5；整机 AME Pure 为左右合并，不乘系数
        df = build_cycles_df(CYCLE_ROWS, "放生料完成,放熟料完成")
        normal = summarize_uph(df, 1, pure_uph_factor=0.5)
        ame = summarize_uph_ame(df, 1)
        self.assertEqual(normal.iloc[0]['Pure UPH(个/小时)'], 180.0)   # 360/2
        self.assertEqual(normal.iloc[0]['Derated UPH M2(个/小时)'], normal.iloc[0]['UPH(个/小时)'])
        self.assertEqual(ame.iloc[0]['Pure UPH(个/小时)'], 360.0)      # 整机不乘系数

    def test_ame_m2_is_sum_of_modules(self):
        # 整机 Derated UPH M2 = 各模组之和（FR 左轴+右轴）
        df = pd.DataFrame({
            "Module": ["左轴", "左轴", "右轴"],
            "Class": ["正常周期"] * 3,
            "CycleSeconds": [3.0, 3.0, 6.0],
            "FileName": ["a.log"] * 3,
        })
        ame = summarize_uph_ame(df, 1)
        self.assertEqual(ame.iloc[0]['Derated UPH M2(个/小时)'], 1200.0 + 600.0)
        self.assertEqual(ame.iloc[0]['Derated UPH M1(个/小时)'], 1200.0 + 600.0)  # 多模组 M1 同样求和

    def test_eff_coretech(self):
        rows = [
            row("2026-07-08 00:00:00.000 [PS][status:RUN,Datetime:...,ReasonID:None]"),
            row("2026-07-08 00:00:10.000 [PS][status:IDLE,Datetime:...,ReasonID:0000010000]"),
            row("2026-07-08 00:00:30.000 [PS][status:DOWN,Datetime:...,ReasonID:0000000411]"),
            row("2026-07-08 00:00:40.000 [PS][status:RUN,Datetime:...,ReasonID:None]"),
            row("2026-07-08 00:01:00.000 [PS][status:IDLE,Datetime:...,ReasonID:0000010000]"),
        ]
        status_summary, _, detail = analyze_status(rows)
        eff = summarize_eff_coretech(status_summary, detail, planned_hours=None, pdt_reason_ids="0000000411")
        r0 = eff.iloc[0]
        self.assertEqual(r0['运行时间RUN(秒)'], 30.0)     # 10 + 20
        self.assertEqual(r0['待机时间IDLE(秒)'], 20.0)
        self.assertEqual(r0['停机时间DOWN(秒)'], 10.0)
        self.assertEqual(r0['计划停机pDT(秒)'], 10.0)     # ReasonID 0000000411 归为计划停机
        self.assertEqual(r0['非计划停机uDT(秒)'], 0.0)
        self.assertEqual(r0['EFF(%)'], 83.33)             # (30+20)/60
        p = down_pareto(detail)
        self.assertEqual(len(p), 1)
        self.assertEqual(p.iloc[0]['ReasonID'], '0000000411')
        self.assertEqual(p.iloc[0]['总时长(秒)'], 10.0)

    def test_reason_code_mapping(self):
        from models.reason_codes import is_planned, load_reason_codes, reason_info
        lm = load_reason_codes("LM")
        self.assertIn("411", lm)
        self.assertEqual(lm["411"]["name"], "打標位建真空")
        self.assertEqual(lm["411"]["category"], "Unplanned Downtime")
        self.assertFalse(is_planned(lm, "411"))
        fr = load_reason_codes("FR")
        self.assertEqual(reason_info(fr, "1110000000")["name"], "真空報警")
        # 计划停机分类：构造 Routine 类型
        fake = {"999": {"name": "例行保养", "category": "Routine Downtime", "state": "DOWN"}}
        self.assertTrue(is_planned(fake, "999"))

    def test_available_reason_lists(self):
        from models.reason_codes import available_reason_lists
        lists = available_reason_lists()
        self.assertIn("LM", lists)
        self.assertIn("FR", lists)

    def test_eff_coretech_with_reason_map(self):
        reason_map = {
            "1110000000": {"name": "真空報警", "category": "Unplanned Downtime", "state": "DOWN"},
            "2220000000": {"name": "例行保養", "category": "Routine Downtime", "state": "DOWN"},
        }
        rows = [
            row("2026-07-08 00:00:00.000 [PS][status:RUN,ReasonID:None]"),
            row("2026-07-08 00:00:10.000 [PS][status:DOWN,ReasonID:1110000000]"),
            row("2026-07-08 00:00:20.000 [PS][status:DOWN,ReasonID:2220000000]"),
            row("2026-07-08 00:00:30.000 [PS][status:RUN,ReasonID:None]"),
        ]
        status_summary, _, detail = analyze_status(rows)
        eff = summarize_eff_coretech(status_summary, detail, reason_map=reason_map)
        self.assertEqual(eff.iloc[0]['非计划停机uDT(秒)'], 10.0)   # 1110000000 → uDT
        self.assertEqual(eff.iloc[0]['计划停机pDT(秒)'], 10.0)     # 2220000000 → pDT
        p = down_pareto(detail, reason_map=reason_map)
        self.assertIn('原因名称', p.columns)
        self.assertIn('停机类型', p.columns)
        self.assertEqual(set(p['停机类型']), {'pDT', 'uDT'})

    def test_eff_pdt_map_plus_manual(self):
        # 有原因清单时，手动计划停机码仍可补充（清单未覆盖的码）
        reason_map = {"1110000000": {"name": "真空報警", "category": "Unplanned Downtime", "state": "DOWN"}}
        rows = [
            row("2026-07-08 00:00:00.000 [PS][status:RUN,ReasonID:None]"),
            row("2026-07-08 00:00:10.000 [PS][status:DOWN,ReasonID:1110000000]"),
            row("2026-07-08 00:00:20.000 [PS][status:DOWN,ReasonID:999]"),
            row("2026-07-08 00:00:30.000 [PS][status:RUN,ReasonID:None]"),
        ]
        status_summary, _, detail = analyze_status(rows)
        eff = summarize_eff_coretech(status_summary, detail, pdt_reason_ids="999", reason_map=reason_map)
        self.assertEqual(eff.iloc[0]['非计划停机uDT(秒)'], 10.0)   # 清单匹配 → uDT
        self.assertEqual(eff.iloc[0]['计划停机pDT(秒)'], 10.0)     # 手动补充 → pDT

    def test_end_to_end_eff(self):
        src = tempfile.mkdtemp()
        out = tempfile.mkdtemp()
        with open(os.path.join(src, "a.log"), "w", encoding="utf-8") as f:
            f.write("2026-07-08 00:00:00.000 [PS][status:RUN,ReasonID:None]\n")
            f.write("2026-07-08 00:00:20.000 [PS][status:IDLE,ReasonID:0000010000]\n")
            f.write("2026-07-08 00:00:40.000 [PS][status:DOWN,ReasonID:0000000411]\n")
            f.write("2026-07-08 00:00:50.000 [PS][status:RUN,ReasonID:None]\n")
        path = LogModel().analyze_eff(src, out, pdt_reason_ids="0000000411")
        xl = pd.ExcelFile(path)
        self.assertIn("Summary", xl.sheet_names)
        self.assertIn("DOWN_Pareto", xl.sheet_names)
        self.assertEqual(xl.parse("Summary").iloc[0]['EFF(%)'], 80.0)  # (20+20)/50

    def test_alarm_boundary(self):
        rows = [
            row("2026-07-08 03:00:00.000   ALARM 123"),
            row("2026-07-08 03:00:05.000   报警:过温"),
            row("2026-07-08 03:00:06.000   ALARM 123"),
            row("2026-07-08 03:00:07.000   无料NG1"),
            row("2026-07-08 03:00:08.000   DNMHVG04XP80000Y2N+1053+B"),
        ]
        summary, by_kw, detail = summarize_alarms(rows, "报警,ALARM,NG")
        self.assertEqual(len(detail), 4)  # 编码中的 ng 不命中
        self.assertEqual(len(by_kw), 3)   # ALARM x2, 报警 x1, NG x1

    def test_alarm_reason_name(self):
        reason_map = {"305": {"name": "吸取钢片真空報警", "category": "Unplanned Downtime", "state": "DOWN"}}
        rows = [
            row("2026-07-08 03:00:00.000 [Error 305] Picker真空报警"),
            row("2026-07-08 03:00:01.000 无料NG"),
        ]
        summary, by_kw, detail = summarize_alarms(rows, "Error,NG", reason_map=reason_map)
        self.assertIn("原因名称", detail.columns)
        self.assertEqual(detail.iloc[0]["原因名称"], "吸取钢片真空報警")
        self.assertEqual(detail.iloc[1]["原因名称"], "")

    def test_status_analysis(self):
        rows = [
            row("2026-07-08 00:00:00.000 Net(2) Send [PS][status:RUN,Datetime:...,ReasonID:None]"),
            row("2026-07-08 00:00:10.000 Net(2) Send [PS][status:IDLE,Datetime:...,ReasonID:0000010000]"),
            row("2026-07-08 00:00:30.000 Net(2) Send [PS][status:RUN,Datetime:...,ReasonID:None]"),
            row("2026-07-08 00:00:50.000 Net(2) Send [PS][status:DOWN,Datetime:...,ReasonID:0000000411]"),
            row("2026-07-08 00:00:56.000 Net(2) Send [PS][status:RUN,Datetime:...,ReasonID:None]"),
        ]
        summary, hourly, detail = analyze_status(rows)
        self.assertEqual(len(detail), 4)
        s = summary.set_index('状态')
        self.assertEqual(s.loc['RUN', '总时长(秒)'], 10.0 + 20.0)  # 最后一条 RUN 无后续事件，不计时长
        self.assertEqual(s.loc['IDLE', '总时长(秒)'], 20.0)
        self.assertEqual(s.loc['DOWN', '总时长(秒)'], 6.0)
        self.assertEqual(len(hourly), 1)

    def test_status_warn_and_rufff(self):
        rows = [
            row("2026-07-08 00:00:00.000 [PS][status:RUFFF,Datetime:...,ReasonID:None]"),
            row("2026-07-08 00:00:10.000 [PS][status:RU100,Datetime:...,ReasonID:None]"),
            row("2026-07-08 00:00:20.000 [PS][status:IDLE,Datetime:...,ReasonID:0000010000]"),
            row("2026-07-08 00:00:30.000 [PS][status:WARN,Datetime:...,ReasonID:0000000411]"),
            row("2026-07-08 00:00:40.000 [PS][status:RUN,Datetime:...,ReasonID:None]"),
        ]
        summary, hourly, detail = analyze_status(rows)
        statuses = set(summary['状态'])
        self.assertIn('RUN', statuses)   # RUFFF 归一化为 RUN
        self.assertIn('WARN', statuses)  # 协议中的警告状态
        self.assertEqual(len(detail), 4)

    def test_em_production(self):
        rows = [
            row("2026-07-28 00:27:39.715 Net(2) Send [startdata][0011901001][000000H6MT][EM]"
                "[lotno:K63030373-24,inputqty:507,goodqty:507,ngqty:0,head:001,"
                "startdatetime:20260728000000491,enddatetime:20260728002739713,datetime:20260728002739715][enddata]"),
        ]
        df = parse_em_production(rows)
        self.assertEqual(len(df), 1)
        r0 = df.iloc[0]
        self.assertEqual(r0['LotNo'], 'K63030373-24')
        self.assertEqual(r0['GoodQty'], 507)
        self.assertEqual(r0['NgQty'], 0)
        duration_h = r0['DurationHours']
        self.assertAlmostEqual(duration_h, (1659.222) / 3600, places=3)
        self.assertAlmostEqual(r0['UPH(个/小时)'], round(507 / (1659.222 / 3600), 2))

    def test_end_to_end_status(self):
        src = tempfile.mkdtemp()
        out = tempfile.mkdtemp()
        with open(os.path.join(src, "a.log"), "w", encoding="utf-8") as f:
            f.write("2026-07-08 00:00:00.000 [status:RUN]\n")
            f.write("2026-07-08 00:01:00.000 [status:IDLE]\n")
        path = LogModel().analyze_status(src, out)
        xl = pd.ExcelFile(path)
        self.assertIn("Summary", xl.sheet_names)
        self.assertIn("Hourly", xl.sheet_names)

    def test_end_to_end_uph(self):
        src = tempfile.mkdtemp()
        out = tempfile.mkdtemp()
        with open(os.path.join(src, "a.log"), "w", encoding="utf-8") as f:
            for c in CYCLE_ROWS:
                f.write(c["Content"] + "\n")
            f.write("2026-07-28 00:27:39.715 [EM][lotno:LOT1,inputqty:100,goodqty:98,ngqty:2,head:001,"
                    "startdatetime:20260728000000491,enddatetime:20260728002739713,datetime:20260728002739715]\n")
        path = LogModel().analyze_uph(src, out, trigger_keywords="放生料完成")
        xl = pd.ExcelFile(path)
        self.assertIn("Summary", xl.sheet_names)
        self.assertIn("AMESummary", xl.sheet_names)
        self.assertIn("CycleDetail", xl.sheet_names)
        self.assertIn("EMProduction", xl.sheet_names)
        self.assertEqual(xl.parse("Summary").iloc[0]['周期总数'], 3)

    def test_analyze_uph_with_file_filters(self):
        src = tempfile.mkdtemp()
        out = tempfile.mkdtemp()
        with open(os.path.join(src, "记录PLC2当前工位当前步数记录.log"), "w", encoding="utf-8") as f:
            for c in CYCLE_ROWS:
                f.write(c["Content"] + "\n")
        with open(os.path.join(src, "RAYPRUS交互记录.log"), "w", encoding="utf-8") as f:
            f.write("2026-07-28 00:27:39.715 [EM][lotno:LOT1,inputqty:100,goodqty:98,ngqty:2,head:001,"
                    "startdatetime:20260728000000491,enddatetime:20260728002739713]\n")
        with open(os.path.join(src, "Debug记录.log"), "w", encoding="utf-8") as f:
            f.write("2026-07-28 00:00:00.000 无料NG\n")
        path = LogModel().analyze_uph(
            src, out, trigger_keywords="放生料完成",
            file_filters=["记录PLC2当前工位当前步数记录", "RAYPRUS交互记录"],
        )
        xl = pd.ExcelFile(path)
        self.assertEqual(xl.parse("Summary").iloc[0]['周期总数'], 3)  # 只读 PLC2 文件
        self.assertEqual(len(xl.parse("EMProduction")), 1)            # EM 来自 RAYPRUS 文件

    def test_custom_template_roundtrip(self):
        from models.process_templates import get_template, load_custom_template, save_custom_template
        tpl = {
            "description": "用户自定义模板",
            "file_filters": ["记录PLC2", "RAYPRUS"],
            "UPH分析": {"trigger_keywords": "放生料完成", "units_per_cycle": 1},
            "EFF分析": {"planned_hours": 22.0, "pdt_reason_ids": "411"},
            "报警分析": {"alarm_keywords": "NG,异常"},
            "机台状态分析": {},
        }
        path = os.path.join(tempfile.mkdtemp(), "custom_template.json")
        save_custom_template(tpl, path)
        loaded = load_custom_template(path)
        self.assertEqual(loaded["UPH分析"]["trigger_keywords"], "放生料完成")
        self.assertEqual(loaded["EFF分析"]["pdt_reason_ids"], "411")
        got = get_template("自定义", path)
        self.assertEqual(got["file_filters"], ["记录PLC2", "RAYPRUS"])

    def test_builtin_templates(self):
        from models.process_templates import PROCESS_TEMPLATES
        self.assertEqual(
            PROCESS_TEMPLATES["LM 激光打标"]["UPH分析"]["trigger_keywords"], "MarkEnd1"
        )
        caw = PROCESS_TEMPLATES["CAW 组装"]
        machines = caw["UPH分析"]["bottleneck_machines"]
        self.assertEqual([m["name"] for m in machines], ["上料机", "焊接机"])
        self.assertEqual(machines[1]["units_per_cycle"], 2)   # 焊接机每周期 2 颗
        self.assertEqual(machines[1]["parallel_units"], 4)    # 左右 2 侧 × 2 焊接头
        self.assertEqual(caw["file_filters"], ["记录PLC", "RAYPRUS", "Debug", "设备状态"])
        self.assertIn("FR 机台", PROCESS_TEMPLATES)
        fr = PROCESS_TEMPLATES["FR 机台"]
        self.assertIn("SA 机台", PROCESS_TEMPLATES)
        self.assertEqual(PROCESS_TEMPLATES["SA 机台"]["UPH分析"]["trigger_keywords"], "Heater 0 :Heating Complete")
        self.assertEqual(PROCESS_TEMPLATES["SA 机台"]["UPH分析"]["units_per_cycle"], 2)
        self.assertEqual(PROCESS_TEMPLATES["SA 机台"]["UPH分析"]["normal_threshold"], 15.0)
        self.assertIn("JigLoadingCycle", PROCESS_TEMPLATES["SA 机台"]["UPH分析"]["tray_change"]["pattern"])
        self.assertEqual(PROCESS_TEMPLATES["SA 机台"]["file_filters"], [".txt"])
        self.assertEqual(PROCESS_TEMPLATES["SA 机台"]["机台状态分析"]["activity_keywords"], "UDP Module - Good")
        self.assertIn("不达标", fr["报警分析"]["alarm_keywords"])
        self.assertIn("有漏点产品", fr["报警分析"]["alarm_keywords"])
        self.assertEqual(fr["UPH分析"]["trigger_keywords"], "轴点胶完成,有漏点产品")
        self.assertEqual(fr["UPH分析"]["module_pattern"], "(左轴|右轴)")
        self.assertEqual(fr["UPH分析"]["pure_uph_factor"], 0.5)

    def test_one_click_excludes_merge(self):
        from controllers.log_controller import ONE_CLICK_FEATURES, FEATURE_METHODS
        # 一键分析只跑 4 项分析，合并保留为独立功能（提速大日志制程）
        self.assertEqual(
            ONE_CLICK_FEATURES,
            ["UPH分析", "EFF分析", "报警分析", "机台状态分析"],
        )
        self.assertNotIn("文档合并与内容拆分", ONE_CLICK_FEATURES)
        self.assertIn("文档合并与内容拆分", FEATURE_METHODS)

    def test_analyze_uph_with_shared_rows(self):
        # 共享已读日志：提供 rows 时不读目录，正常出结果
        rows = [{"FileName": c["FileName"], "Content": c["Content"]} for c in CYCLE_ROWS]
        out = tempfile.mkdtemp()
        path = LogModel().analyze_uph("不存在的目录", out, trigger_keywords="放生料完成", rows=rows)
        self.assertEqual(pd.ExcelFile(path).parse("Summary").iloc[0]['周期总数'], 3)

    def test_cancel_event_skips_cycle_build(self):
        import threading
        from models.exceptions import OperationCancelled
        cancel = threading.Event()
        cancel.set()
        with self.assertRaises(OperationCancelled):
            build_cycles_df(CYCLE_ROWS, "放生料完成", cancel_event=cancel)

    def test_cancel_event_skips_read_files(self):
        import threading
        from models.exceptions import OperationCancelled
        from utils.file_utils import read_files
        src = tempfile.mkdtemp()
        with open(os.path.join(src, "a.log"), "w", encoding="utf-8") as f:
            f.write("x\n")
        cancel = threading.Event()
        cancel.set()
        with self.assertRaises(OperationCancelled):
            read_files(src, cancel_event=cancel)

    def test_module_pattern_with_defect_completion(self):
        # 漏点件以"有漏点产品"代替"点胶完成"，周期序列应连续且分左/右轴
        rows = [
            {"FileName": "a.log", "Content": "[00:00:00] 左轴点胶完成"},
            {"FileName": "a.log", "Content": "[00:00:06] 右轴点胶完成"},
            {"FileName": "a.log", "Content": "[00:00:07] 左轴有漏点产品，请处理！"},
            {"FileName": "a.log", "Content": "[00:00:14] 右轴有漏点产品，请处理！"},
            {"FileName": "a.log", "Content": "[00:00:15] 左轴点胶完成"},
        ]
        df = build_cycles_df(rows, "轴点胶完成,有漏点产品",
                             module_pattern=r"(左轴|右轴)")
        self.assertEqual(set(df["Module"]), {"左轴", "右轴"})
        left = df[df["Module"] == "左轴"]
        right = df[df["Module"] == "右轴"]
        self.assertEqual(list(left["CycleSeconds"]), [7.0, 8.0])   # 漏点件计入周期
        self.assertEqual(list(right["CycleSeconds"]), [8.0])

    def test_ppt_report(self):
        from models.report import build_ppt_report
        from pptx import Presentation
        src = tempfile.mkdtemp()
        out = tempfile.mkdtemp()
        lines = [
            "2026-07-08 00:00:00.000 [PS][status:RUN,ReasonID:None]",
            "2026-07-08 00:00:10.000 [PS][status:IDLE,ReasonID:0000010000]",
            "2026-07-08 00:00:30.000 [PS][status:DOWN,ReasonID:0000000411]",
            "2026-07-08 00:00:40.000 [PS][status:RUN,ReasonID:None]",
            "2026-07-08 00:00:50.000 无料NG",
            "2026-07-08 00:00:55.000 MarkEnd1",
            "2026-07-08 00:00:58.000 MarkEnd1",
            "2026-07-08 00:00:02.000 [EM][lotno:LOT1,inputqty:100,goodqty:98,ngqty:2,head:001,"
            "startdatetime:20260708000000491,enddatetime:20260708000030000]",
        ]
        with open(os.path.join(src, "a.log"), "w", encoding="utf-8") as f:
            f.write("\n".join(lines) + "\n")
        m = LogModel()
        m.analyze_uph(src, out, trigger_keywords="MarkEnd1")
        m.analyze_eff(src, out)
        m.analyze_alarms(src, out)
        m.analyze_status(src, out)
        ppt = build_ppt_report(out, process_name="CAW 组装")
        prs = Presentation(ppt)
        self.assertEqual(len(prs.slides), 7)  # 封面 + 6 内容页（模板目录页已删除）
        self.assertFalse(any('agenda' in s.slide_layout.name.lower() for s in prs.slides))
        self.assertTrue(any(shape.has_chart for s in prs.slides for shape in s.shapes))
        cover_text = "\n".join(
            sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
        )
        self.assertIn("CAW設備一鍵自動分析報告", cover_text)  # 制程字母已替换
        self.assertIn("——UPH、EFF、Alarm", cover_text)  # 模板副标题保留
        self.assertIn("1 /", cover_text)
        self.assertIn("7", cover_text)  # 封面总页数
        uph_page = "\n".join(
            sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame
        )
        self.assertIn("2/7", uph_page)  # 内容页页码同步（与模板样式一致）
        # 原因码前导零在 PPT 表格中保留（0000000411）
        table_texts = []
        for s in prs.slides:
            for shape in s.shapes:
                if shape.has_table:
                    for r in range(len(shape.table.rows)):
                        for c in range(len(shape.table.columns)):
                            table_texts.append(shape.table.cell(r, c).text)
        self.assertIn("0000000411", table_texts)

    def test_ppt_report_nan_values(self):
        # 回归：AMESummary 存在 NaN（如 FR 无周期时 Pure/M2 为空）不应导致图表写入报错
        from models.report import build_ppt_report
        from pptx import Presentation
        out = tempfile.mkdtemp()
        LogModel._write_sheets(os.path.join(out, "UPH_Analysis.xlsx"), {
            "AMESummary": pd.DataFrame({
                "Pure UPH(个/小时)": [float("nan")],
                "Derated UPH M1(个/小时)": [2892.32],
                "Derated UPH M2(个/小时)": [float("nan")],
            })
        })
        ppt = build_ppt_report(out, process_name="FR 机台")
        self.assertTrue(os.path.exists(ppt))
        prs = Presentation(ppt)
        self.assertGreaterEqual(len(prs.slides), 2)  # 封面 + UPH 页（其余 Excel 缺失被移除）
        self.assertTrue(any(shape.has_chart for s in prs.slides for shape in s.shapes))

    def test_ppt_report_bottleneck_mode(self):
        from models.report import build_ppt_report
        from pptx import Presentation
        out = tempfile.mkdtemp()
        LogModel._write_sheets(os.path.join(out, "UPH_Analysis.xlsx"), {
            "Summary": pd.DataFrame({
                "工位": ["点胶", "贴附", "热压", "检测"],
                "事件数": [1456, 2825, 1414, 11061],
                "每排事件数": [1, 2, 1, 8],
                "每排周期(秒)": [9.964, 9.899, 11.061, 9.961],
                "极限UPH(个/小时)": [722.6, 727.35, 650.94, 722.82],
                "瓶颈": ["", "", "★", ""],
            }),
            "AMESummary": pd.DataFrame({
                "瓶颈工位": ["热压"],
                "瓶颈周期(秒)": [11.061],
                "每排产品数(个)": [2],
                "Pure UPH(个/小时)": [650.94],
            }),
        })
        ppt = build_ppt_report(out, process_name="SA 机台")
        prs = Presentation(ppt)
        uph_slide = next(
            s for s in prs.slides
            if s.shapes.title is not None and s.shapes.title.text == "UPH 分析"
        )
        texts = "\n".join(sh.text_frame.text for sh in uph_slide.shapes if sh.has_text_frame)
        self.assertIn("瓶颈工位：热压", texts)
        self.assertIn("650.94", texts)
        # PPT 只放最终结果：不展示工位明细表
        self.assertFalse(any(sh.has_table for sh in uph_slide.shapes))
        self.assertTrue(any(sh.has_chart for sh in uph_slide.shapes))

    def test_status_time_only_multi_file(self):
        # FR 风格：纯时间戳 + 按日分文件，跨日 23:59:58 -> 00:00:01 应为 3 秒
        rows = [
            {"FileName": "L06-25.log", "Content": "[23:59:58] Send: [PS][status:RUN,ReasonID:None]"},
            {"FileName": "L06-26.log", "Content": "[00:00:01] Send: [PS][status:IDLE,ReasonID:0020000000]"},
            {"FileName": "L06-26.log", "Content": "[00:00:31] Send: [PS][status:RUN,ReasonID:None]"},
        ]
        summary, hourly, detail = analyze_status(rows)
        self.assertEqual(len(detail), 2)
        s = summary.set_index('状态')
        self.assertEqual(s.loc['RUN', '总时长(秒)'], 3.0)   # 跨日 3 秒；最后一条 RUN 无后续事件不计时长
        self.assertEqual(s.loc['IDLE', '总时长(秒)'], 30.0)

    def test_status_derived(self):
        reason_map = {"305": {"name": "吸取钢片真空報警", "category": "Unplanned Downtime", "state": "DOWN"}}
        rows = [
            {"FileName": "a.txt", "Content": "[0] X, 20260803-09-00-00-000, X, 0, UDP Module - Good"},
            {"FileName": "a.txt", "Content": "[1] X, 20260803-09-00-05-000, X, 0, UDP Module - Good"},
            {"FileName": "a.txt", "Content": "[2] X, 20260803-09-00-08-000, X, 0, AutoRun Stop - ErrorName = [Error 305] Picker"},
            {"FileName": "a.txt", "Content": "[3] X, 20260803-09-00-12-000, X, 0, UDP Module - Good"},
            {"FileName": "a.txt", "Content": "[4] X, 20260803-09-00-15-000, X, 0, AutoRun Stop - ErrorName = [작업자 정지]"},
            {"FileName": "a.txt", "Content": "[5] X, 20260803-09-00-20-000, X, 0, UDP Module - Good"},
        ]
        summary, hourly, detail = analyze_status_derived(
            rows, "UDP Module - Good", "AutoRun Stop - ErrorName", reason_map=reason_map,
        )
        s = summary.set_index('状态')
        self.assertEqual(s.loc['RUN', '总时长(秒)'], 8.0 + 3.0)   # 0→8、12→15
        self.assertEqual(s.loc['DOWN', '总时长(秒)'], 4.0)         # Error 305 → DOWN
        self.assertEqual(s.loc['IDLE', '总时长(秒)'], 5.0)         # 操作员停止 → IDLE
        down = detail[detail['Status'] == 'DOWN'].iloc[0]
        self.assertEqual(down['ReasonID'], '305')

    def test_analyze_bottleneck(self):
        import datetime as _dt
        from models.analysis import analyze_bottleneck

        def line(msg, t):
            s = t.strftime('%Y%m%d-%H-%M-%S-') + f"{t.microsecond // 1000:03d}"
            return {"FileName": "a.txt", "Content": f"[0] X, {s}, X, 0, {msg}"}

        t0 = _dt.datetime(2026, 8, 3, 9, 0, 0)
        rows = []
        for i in range(3):   # 点胶：每排 5s
            rows.append(line("DispOneChipProfileWorkCycle", t0 + _dt.timedelta(seconds=i * 5)))
        for i in range(6):   # 贴附：每排 2 次，每排 5s
            rows.append(line("AfterPickUp StopCondition", t0 + _dt.timedelta(seconds=i * 2.5)))
        for i in range(3):   # 热压：每排 10s（瓶颈）
            rows.append(line("Heater 0 :Heating Complete", t0 + _dt.timedelta(seconds=i * 10)))
        for i in range(6):   # 检测：每排 2 次，每排 2s
            rows.append(line("UDP Module - Good", t0 + _dt.timedelta(seconds=i * 1)))

        stations = [
            {"name": "点胶", "function": "sa_dispense"},
            {"name": "贴附", "function": "sa_attach"},
            {"name": "热压", "function": "sa_heatpress"},
            {"name": "检测", "function": "sa_inspect"},
        ]
        df, result = analyze_bottleneck(rows, stations, units_per_row=2)
        self.assertEqual(result["瓶颈工位"], "热压")
        self.assertAlmostEqual(result["瓶颈周期(秒)"], 10.0)
        self.assertEqual(df.set_index("工位").loc["热压", "瓶颈"], "★")

    def test_analyze_bottleneck_tray_amortize(self):
        import datetime as _dt
        from models.analysis import analyze_bottleneck

        def line(msg, t):
            s = t.strftime('%Y%m%d-%H-%M-%S-') + f"{t.microsecond // 1000:03d}"
            return {"FileName": "a.txt", "Content": f"[0] X, {s}, X, 0, {msg}"}

        t0 = _dt.datetime(2026, 8, 3, 9, 0, 0)
        rows = []
        for t in (0, 20, 25):   # 热压
            rows.append(line("Heater 0 :Heating Complete", t0 + _dt.timedelta(seconds=t)))
        for t in (0, 10, 15):   # 点胶
            rows.append(line("DispOneChipProfileWorkCycle", t0 + _dt.timedelta(seconds=t)))
        rows.append(line("JigUnloadingCycle", t0 + _dt.timedelta(seconds=5)))   # 卸盘
        rows.append(line("JigLoadingCycle", t0 + _dt.timedelta(seconds=8)))     # 装盘

        stations = [
            {"name": "点胶", "function": "sa_dispense"},
            {"name": "热压", "function": "sa_heatpress"},
        ]
        df, result = analyze_bottleneck(rows, stations, units_per_row=2,
                                        tray_change={
                                            "pattern": "JigLoadingCycle",
                                            "unload_pattern": "JigUnloadingCycle",
                                            "rows_per_tray": 2,
                                            "single_tray_seconds": 10,
                                        })
        self.assertEqual(result["瓶颈工位"], "热压")
        self.assertAlmostEqual(result["瓶颈周期(秒)"], 12.5)      # median(20, 5)
        self.assertAlmostEqual(result["单次换盘时间(秒)"], 10.0)   # 配置覆盖
        self.assertAlmostEqual(result["每排换盘开销(秒)"], 5.0)   # 10s / 2排
        self.assertAlmostEqual(result["有效周期(秒)"], 17.5)      # 12.5 + 5

    def test_analyze_uph_parts(self):
        import datetime as _dt

        def line(fname, msg, t):
            s = t.strftime('%Y%m%d-%H-%M-%S-') + f"{t.microsecond // 1000:03d}"
            return {"FileName": fname, "Content": f"[0] X, {s}, X, 0, {msg}"}

        t0 = _dt.datetime(2026, 8, 5, 0, 0, 0)
        rows = []
        for i in range(3):
            rows.append(line("上料機/Logger.txt", "吸嘴放料成功", t0 + _dt.timedelta(seconds=i * 4)))
        for i in range(3):
            rows.append(line("主機/OHLog.txt", "ProcessJob,Cavity cnt:1,K1,OK", t0 + _dt.timedelta(seconds=i * 5)))
        for i in range(3):
            rows.append(line("下料機/Logger.txt", "UnloadDuts Finish", t0 + _dt.timedelta(seconds=i * 6)))
        parts = [
            {"name": "上料機", "trigger": "放料成功", "units_per_cycle": 1},
            {"name": "主機", "trigger": "Cavity cnt:", "units_per_cycle": 1},
            {"name": "下料機", "trigger": "UnloadDuts Finish", "units_per_cycle": 1},
        ]
        out = tempfile.mkdtemp()
        path = LogModel().analyze_uph("不存在的目录", out, parts=parts, module_from_path=True, rows=rows)
        df = pd.ExcelFile(path).parse("Summary")
        self.assertEqual(set(df["模块"]), {"上料機", "主機", "下料機"})
        self.assertEqual(df.set_index("模块").loc["上料機", "周期总数"], 2)
        self.assertEqual(df.set_index("模块").loc["主機", "周期总数"], 2)
        self.assertEqual(df.set_index("模块").loc["下料機", "周期总数"], 2)

    def test_alarms_module_from_path(self):
        rows = [
            {"FileName": "上料機/a.txt", "Content": "2026-07-08 00:00:00.000 NG abc"},
            {"FileName": "下料機/b.txt", "Content": "2026-07-08 00:00:01.000 NG xyz"},
        ]
        summary, by_kw, detail = summarize_alarms(rows, "NG", module_from_path=True)
        self.assertEqual(set(summary["模块"]), {"上料機", "下料機"})


if __name__ == "__main__":
    unittest.main()
