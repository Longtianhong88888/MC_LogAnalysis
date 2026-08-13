"""联网一键报告（PPT）测试。"""

import os
import tempfile
import unittest
from unittest.mock import patch

import pandas as pd
from pptx import Presentation

from models.analysis import build_web_uph_sheets, summarize_web_alarms
from models import web_api
from models.log_model import LogModel


def runlog_row(**kw):
    base = {
        "happentime": "2026-08-12 08:00:00",
        "endtime": "2026-08-12 08:18:57",
        "status": "DOWN",
        "reasonid": "0000010792",
        "errorname": "Routine cleaning machine",
        "errorMsg": "常規清潔機台",
        "downFlag": "Routine downtime",
        "spendTime": 18.95,
    }
    base.update(kw)
    return base


class WebSheetsTest(unittest.TestCase):
    def test_build_web_uph_sheets(self):
        output_rows = [
            {"machine_no": "CAW7203", "machine_type": "CAW", "device_name": "AKC",
             "status": "RUN", "input": 11036, "output": 10940, "output_target": 29000,
             "hit_rate": 0.3773},
            {"machine_no": "CAW7205", "machine_type": "CAW", "device_name": "ATW-E",
             "status": "IDLE", "input": 2, "output": 1, "output_target": 0,
             "hit_rate": 0},
        ]
        eff_rows = [{"machineNO": "CAW7203", "uph": "982"}]
        summary, ame = build_web_uph_sheets(
            output_rows, eff_rows,
            begin_time="2026/08/12 06:00:00", end_time="2026/08/13 06:00:00",
        )
        self.assertEqual(list(summary["机台号"]), ["CAW7203", "CAW7205"])
        self.assertEqual(summary.iloc[0]["产出(个)"], 10940)
        self.assertEqual(summary.iloc[0]["UPH(个/小时)"], 982)   # eff 兜底
        self.assertAlmostEqual(summary.iloc[1]["UPH(个/小时)"], 0.04, delta=0.001)  # 1/24 保留两位
        self.assertEqual(ame.iloc[0]["总产出(个)"], 10941)
        self.assertEqual(ame.iloc[0]["机台数"], 2)
        self.assertEqual(ame.iloc[0]["数据来源"], "联网接口(CMS)")

    def test_summarize_web_alarms(self):
        raw_logs = [
            {"machineNo": "CAW7203", "rows": [
                runlog_row(),
                runlog_row(happentime="2026-08-12 08:33:17", endtime="2026-08-12 08:33:58",
                           status="DOWN", reasonid="0000010629",
                           errorname="Alam[16,17]^Chassis Left vacuum",
                           errorMsg="Alam[16,17]^Chassis左取料模組吸真空1",
                           downFlag="Unplanned Downtime"),
                runlog_row(happentime="2026-08-12 08:00:00", endtime="2026-08-12 08:10:00",
                           status="RUN", downFlag=""),
            ]},
            {"machineNo": "CAW7205", "rows": [runlog_row()]},
        ]
        summary, by_kw, detail = summarize_web_alarms(raw_logs)
        self.assertEqual(summary["报警次数"].sum(), 3)          # 2 + 1 个 DOWN 段
        self.assertEqual(len(detail), 3)
        self.assertEqual(summary.iloc[0]["模块"], "CAW7203")
        self.assertEqual(detail.iloc[0]["原因名称"], "常規清潔機台")
        # 中文报错优先（英文 errorname 常超长，图表/表格易截断）
        self.assertIn("常規清潔機台", by_kw["命中关键词"].values)
        self.assertIn("Alam[16,17]^Chassis左取料模組吸真空1", by_kw["命中关键词"].values)
        # 报警次数按降序排列
        self.assertEqual(
            list(by_kw["报警次数"]),
            sorted(by_kw["报警次数"], reverse=True),
        )


class AnalyzeWebReportTest(unittest.TestCase):
    MACHINES = [
        {"machineNo": "CAW7203", "machine": "CAW"},
        {"machineNo": "CAW7205", "machine": "CAW"},
        {"machineNo": "LM4401", "machine": "LM"},
    ]
    EFF_ROWS = [
        {"machineNO": "CAW7203", "machineType": "CAW", "status": "RUN", "uph": "982",
         "runTime": "86.62", "idleTime": "10.83", "errorTime": "1342.55",
         "plannedDT": "100", "unPlannedDT": "1242.55", "eff": "6.77%",
         "outputList": [{"data": "2026/08/12", "value": 10940}]},
        {"machineNO": "CAW7205", "machineType": "CAW", "status": "IDLE", "uph": "0",
         "runTime": "0", "idleTime": "1440", "errorTime": "0",
         "plannedDT": "0", "unPlannedDT": "0", "eff": "100.00%",
         "outputList": [{"data": "2026/08/12", "value": 1}]},
    ]
    OUTPUT_ROWS = [
        {"machine_no": "CAW7203", "machine_type": "CAW", "device_name": "AKC",
         "status": "RUN", "input": 11036, "output": 10940, "output_target": 29000,
         "hit_rate": 0.3773},
        {"machine_no": "CAW7205", "machine_type": "CAW", "device_name": "ATW-E",
         "status": "IDLE", "input": 2, "output": 1, "output_target": 0, "hit_rate": 0},
    ]
    RUNLOG = [
        runlog_row(),
        runlog_row(happentime="2026-08-12 08:18:57", endtime="2026-08-12 08:33:17",
                   status="RUN", errorname="", errorMsg="", downFlag="", spendTime=14.33),
        runlog_row(happentime="2026-08-12 08:33:17", endtime="2026-08-12 08:33:58",
                   status="DOWN", reasonid="0000010629",
                   errorname="Alam[16,17]^Chassis Left vacuum",
                   errorMsg="Alam[16,17]^Chassis左取料模組吸真空1",
                   downFlag="Unplanned Downtime", spendTime=0.6833),
    ]

    def test_end_to_end(self):
        out = tempfile.mkdtemp()

        def fake_machines(base_url, plant_id, **kw):
            return self.MACHINES

        def fake_eff(base_url, plant_id, machine_type="", **kw):
            return self.EFF_ROWS if machine_type == "CAW" else []

        def fake_output(base_url, plant_id, machine_type, **kw):
            return self.OUTPUT_ROWS if machine_type == "CAW" else []

        def fake_runlog(base_url, plant_id, machine_no, **kw):
            return self.RUNLOG if machine_no == "CAW7203" else []

        with patch("models.web_api.fetch_machines", side_effect=fake_machines), \
             patch("models.web_api.fetch_machine_eff", side_effect=fake_eff), \
             patch("models.web_api.fetch_machine_output_board", side_effect=fake_output), \
             patch("models.web_api.fetch_run_log", side_effect=fake_runlog):
            result = LogModel().analyze_web_report(
                output_dir=out, process_name="CAW 组装",
                api_url="http://10.151.128.35:8098", plant_id="8S01",
                machine_type="CAW", machine_nos="",
                begin_time="2026/08/12 06:00:00", end_time="2026/08/13 06:00:00",
            )

        ppt_path = result.split("（")[0]
        self.assertTrue(os.path.exists(ppt_path))
        self.assertTrue(os.path.exists(os.path.join(out, "UPH_Analysis.xlsx")))
        self.assertTrue(os.path.exists(os.path.join(out, "EFF_Analysis.xlsx")))
        self.assertTrue(os.path.exists(os.path.join(out, "Alarm_Analysis.xlsx")))
        self.assertTrue(os.path.exists(os.path.join(out, "Status_Analysis.xlsx")))

        uph = pd.ExcelFile(os.path.join(out, "UPH_Analysis.xlsx"))
        self.assertIn("Summary", uph.sheet_names)
        self.assertIn("AMESummary", uph.sheet_names)
        ame = uph.parse("AMESummary")
        self.assertEqual(ame.iloc[0]["数据来源"], "联网接口(CMS)")
        self.assertEqual(ame.iloc[0]["总产出(个)"], 10941)
        s = uph.parse("Summary")
        self.assertEqual(set(s["机台号"]), {"CAW7203", "CAW7205"})

        eff = pd.ExcelFile(os.path.join(out, "EFF_Analysis.xlsx"))
        self.assertIn("DOWN_Pareto", eff.sheet_names)
        eff_summary = eff.parse("Summary").iloc[0]
        self.assertEqual(eff_summary["计划停机pDT(秒)"], 1137.0)
        self.assertEqual(eff_summary["非计划停机uDT(秒)"], 41.0)

        alarm = pd.ExcelFile(os.path.join(out, "Alarm_Analysis.xlsx"))
        self.assertIn("Summary", alarm.sheet_names)
        self.assertIn("ByKeyword", alarm.sheet_names)
        self.assertIn("Detail", alarm.sheet_names)
        self.assertEqual(alarm.parse("Summary")["报警次数"].sum(), 2)

        status = pd.ExcelFile(os.path.join(out, "Status_Analysis.xlsx"))
        self.assertIn("Hourly", status.sheet_names)
        self.assertIn("Detail", status.sheet_names)

        prs = Presentation(ppt_path)
        self.assertGreaterEqual(len(prs.slides), 5)  # 封面 + UPH/EFF/报警/状态
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        joined = "\n".join(texts)
        self.assertIn("总产出 10941", joined)
        # 报告逻辑流：趋势页前置到第 2 页（封面后）
        slide2_texts = "\n".join(
            sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame
        )
        self.assertIn("机台状态趋势", slide2_texts)
        # 封面 KPI 摘要卡 + 页头 chip 联网口径
        cover_texts = "\n".join(
            sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
        )
        self.assertIn("综合EFF", cover_texts)
        self.assertIn("总产出", cover_texts)
        self.assertIn("设备能效", joined)
        # 封面站位前缀来自页面选的站位（而非制程模板）
        self.assertIn("CAW設備能效報告", cover_texts)
        # 停机 Pareto 中文原因名称
        eff = pd.ExcelFile(os.path.join(out, "EFF_Analysis.xlsx"))
        pareto = eff.parse("DOWN_Pareto")
        self.assertIn("原因名称", pareto.columns)
        self.assertNotEqual(str(pareto.iloc[0]["原因名称"]).strip(), "")

    def test_fetch_preview_and_no_data(self):
        out = tempfile.mkdtemp()

        def fake_machines(base_url, plant_id, **kw):
            return self.MACHINES

        def fake_eff(base_url, plant_id, machine_type="", **kw):
            return self.EFF_ROWS if machine_type == "CAW" else []

        def fake_output(base_url, plant_id, machine_type, **kw):
            return self.OUTPUT_ROWS if machine_type == "CAW" else []

        def fake_runlog(base_url, plant_id, machine_no, **kw):
            return self.RUNLOG if machine_no == "CAW7203" else []

        with patch("models.web_api.fetch_machines", side_effect=fake_machines), \
             patch("models.web_api.fetch_machine_eff", side_effect=fake_eff), \
             patch("models.web_api.fetch_machine_output_board", side_effect=fake_output), \
             patch("models.web_api.fetch_run_log", side_effect=fake_runlog):
            data = LogModel().fetch_web_report_data(
                api_url="http://10.151.128.35:8098", plant_id="8S01",
                machine_type="CAW", machine_nos="",
                begin_time="2026/08/12 06:00:00", end_time="2026/08/13 06:00:00",
            )
        self.assertIn("查询成功：2 台机台", data["preview_text"])
        self.assertIn("CAW7203", data["preview_text"])
        self.assertIn("数据就绪", data["preview_text"])

        # 无数据 → WebNoDataError
        def empty_runlog(base_url, plant_id, machine_no, **kw):
            return []

        with patch("models.web_api.fetch_machines", side_effect=fake_machines), \
             patch("models.web_api.fetch_machine_eff", side_effect=lambda *a, **k: []), \
             patch("models.web_api.fetch_machine_output_board", side_effect=lambda *a, **k: []), \
             patch("models.web_api.fetch_run_log", side_effect=empty_runlog):
            with self.assertRaises(web_api.WebNoDataError):
                LogModel().fetch_web_report_data(
                    api_url="http://10.151.128.35:8098", plant_id="8S01",
                    machine_type="CAW", machine_nos="",
                    begin_time="2026/08/12 06:00:00", end_time="2026/08/13 06:00:00",
                )

    def test_report_filename(self):
        self.assertEqual(
            LogModel._web_report_filename({
                "query": {"machine_type": "CAW", "machine_nos": "CAW7203,CAW7205"},
                "begin_time": "2026/08/12 06:00:00",
                "end_time": "2026/08/13 06:00:00",
            }),
            "CAW_CAW7203_CAW7205_20260812-20260813.pptx",
        )
        self.assertEqual(
            LogModel._web_report_filename({
                "query": {"machine_type": "FR", "machine_nos": ""},
                "begin_time": "2026-07-02 06:00:00",
                "end_time": "2026-07-03 06:00:00",
            }),
            "FR_ALL_20260702-20260703.pptx",
        )


if __name__ == "__main__":
    unittest.main()
