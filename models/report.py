"""PPT 报告生成：读取一键分析产出的 4 个 Excel，基于最终模板 Analysis_Report.pptx 生成报告。

模板封面保留并替换制程字母；内容页原位复用（更新副标题/图表/表格），
页码按 '1 /' + 总页数（封面）与 'n/总页数'（内容页）同步更新。
"""

import copy
import io
import math
import os
import re
from datetime import datetime
from collections import Counter

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from models.analysis import fmt_duration
from utils.resource_utils import find_external_resource

HEADER_FILL = RGBColor(0xDD, 0xEB, 0xF7)
GRAY = RGBColor(0x60, 0x60, 0x60)
CHART_LEFT, CHART_TOP, CHART_W, CHART_H = 0.6, 1.6, 6.6, 5.2
TABLE_LEFT, TABLE_W = 7.4, 5.3
_TITLE_PREFIX = {
    "LM 激光打标": "LM",
    "CAW 组装": "CAW",
    "FR 机台": "FR",
    "SA 机台": "SA",
    "ACF 三機": "ACF",
}
_TEMPLATE_CANDIDATES = ["Analysis_Report.pptx", "PPT模板.pptx"]
SECTION_TITLES = ["UPH 分析", "EFF 分析", "停机 Pareto", "报警分析", "机台状态汇总", "机台状态趋势"]

# 方案C 成品视觉稿（Analysis_Report.pptx）设计参数：
# 600×338 设计单位画布 → 13.333×7.5in，表格模块位于 350..574 × 86..300
_TEMPLATE_DESIGN = {
    "table_box": (350, 86, 574, 300),   # 表格模块（设计单位）
    "table_font_pt": 13,                # 表格字号
    "table_min_row_in": 0.25,           # 行高下限（in），13pt 文本安全高度
    "table_max_data_rows": 18,          # 数据行上限（超出截断，防止溢出）
    "theme_bg": "0D1524",               # 深色驾驶舱背景
    "theme_panel": "131E33",            # 模块面板
    "theme_panel_line": "25344F",       # 面板描边
    "theme_title": "E7F1FF",            # 主标题
    "theme_sub": "8FA3C4",              # 副标题
    "theme_accent": "F5B64C",           # 强调色（琥珀）
    "theme_cyan": "35C6F4",             # 图表主色（青）
}


def _read_sheet(path, name):
    if not os.path.exists(path):
        return pd.DataFrame()
    try:
        # dtype=str：保留原因码等前导零（如 0010000000），图表取值再转数值
        return pd.read_excel(path, sheet_name=name, dtype=str)
    except Exception:
        return pd.DataFrame()


def _to_float(value):
    try:
        if value is None or (isinstance(value, str) and value.strip() == ''):
            return None
        f = float(value)
        if f != f or f in (float('inf'), float('-inf')):
            return None
        return f
    except (TypeError, ValueError):
        return None


def _chart_values(values):
    """图表数值清洗：NaN/Inf/空串统一转 0，避免 python-pptx 写入报错。"""
    out = []
    for value in values:
        try:
            f = float(value)
        except (TypeError, ValueError):
            f = None
        if f is None or f != f or f in (float('inf'), float('-inf')):
            out.append(0)
        else:
            out.append(f)
    return out


def _set_text_preserving(tf, text):
    """替换文本框内容并保留首 run 的样式。"""
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = str(text)
        for extra in list(tf.paragraphs[0].runs[1:]):
            extra._r.getparent().remove(extra._r)
        for para in list(tf.paragraphs[1:]):
            para._p.getparent().remove(para._p)
    else:
        tf.text = str(text)


def _content_layout(prs):
    for layout in prs.slide_layouts:
        if "content" in layout.name.lower():
            return layout
    return prs.slide_layouts[0] if prs.slide_layouts else None


def _clean_placeholders(slide):
    """删除内容页上空白的 body/clipart 占位符（保留标题、页脚、页码）。"""
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        idx = shape.placeholder_format.idx
        if idx == 0 or idx in (3, 4):
            continue
        sp = shape._element
        sp.getparent().remove(sp)


def _add_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(_content_layout(prs))
    _clean_placeholders(slide)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8))
        box.text_frame.text = title
        box.text_frame.paragraphs[0].font.size = Pt(26)
        box.text_frame.paragraphs[0].font.bold = True
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.18), Inches(12.3), Inches(0.32))
        tf2 = box2.text_frame
        tf2.text = subtitle
        tf2.paragraphs[0].font.size = Pt(13)
        tf2.paragraphs[0].font.color.rgb = GRAY
    return slide


def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[-1] if prs.slide_layouts else None


def _add_dark_gantt_slide(prs, index, total):
    """新增甘特页：与深色驾驶舱模板同一设计系统（深底/标签/胶囊/面板/页脚）。"""
    slide = prs.slides.add_slide(_blank_layout(prs))
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0D, 0x15, 0x24)
    bg.line.fill.background()
    bg.shadow.inherit = False

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.53), Inches(0.33), Inches(1.6), Inches(0.4))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(0x16, 0x23, 0x3C)
    tag.line.fill.background()
    tag.shadow.inherit = False
    tag.text_frame.text = "机台日志分析"
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag.text_frame.paragraphs[0].font.size = Pt(13)
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x35, 0xC6, 0xF4)

    tb = slide.shapes.add_textbox(Inches(2.67), Inches(0.27), Inches(8.5), Inches(0.8))
    tb.text_frame.text = "瓶颈工序甘特图"
    tb.text_frame.paragraphs[0].font.size = Pt(40)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xE7, 0xF1, 0xFF)

    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(11.4), Inches(0.36), Inches(1.93), Inches(0.4))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(0x16, 0x23, 0x3C)
    chip.line.fill.background()
    chip.shadow.inherit = False
    chip.text_frame.text = f"{index:02d} / {total:02d}"
    chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    chip.text_frame.paragraphs[0].font.size = Pt(13)
    chip.text_frame.paragraphs[0].font.bold = True
    chip.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x35, 0xC6, 0xF4)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.4), Inches(1.78), Inches(12.53), Inches(4.75))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x13, 0x1E, 0x33)
    panel.line.color.rgb = RGBColor(0x25, 0x33, 0x4F)
    panel.line.width = Pt(0.75)
    panel.shadow.inherit = False
    try:
        panel.adjustments[0] = 0.03
    except Exception:
        pass
    for x0, y0, dx, dy in (
        (0.4, 1.78, 1, 1), (12.93, 1.78, -1, 1),
        (0.4, 6.53, 1, -1), (12.93, 6.53, -1, -1),
    ):
        for x1, y1 in ((x0 + dx * 0.18, y0), (x0, y0 + dy * 0.18)):
            c = slide.shapes.add_connector(1, Inches(x0), Inches(y0),
                                           Inches(x1), Inches(y1))
            c.line.color.rgb = RGBColor(0x35, 0xC6, 0xF4)
            c.line.width = Pt(1)
            c.shadow.inherit = False

    line = slide.shapes.add_connector(1, Inches(0.53), Inches(6.93),
                                      Inches(12.8), Inches(6.93))
    line.line.color.rgb = RGBColor(0x22, 0x31, 0x4D)
    line.line.width = Pt(0.75)
    line.shadow.inherit = False

    f = slide.shapes.add_textbox(Inches(0.53), Inches(7.13), Inches(8), Inches(0.3))
    f.text_frame.text = "© 2026 Hon Hai Precision Industry Co., Ltd. All rights reserved."
    f.text_frame.paragraphs[0].font.size = Pt(13)
    f.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x5F, 0x73, 0x96)

    pn = slide.shapes.add_textbox(Inches(11.56), Inches(7.13), Inches(1.24), Inches(0.3))
    pn.text_frame.text = f"{index}/{total}"
    pn.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    pn.text_frame.paragraphs[0].font.size = Pt(14)
    pn.text_frame.paragraphs[0].font.bold = True
    pn.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xF5, 0xB6, 0x4C)
    return slide


def _display_value(col, value):
    """PPT 表格单元格显示：时间列（列名含 秒/时长）按 <60s 秒、>=60s h:mm:ss 展示；
    图表仍取原始数值，不受影响。"""
    if pd.isna(value):
        return ''
    if isinstance(col, str) and ('秒' in col or '时长' in col):
        try:
            return fmt_duration(float(value), with_ms=False)
        except (TypeError, ValueError):
            return str(value)
    return str(value)


def _add_table(slide, df, left, top, width, height, max_rows=12):
    data = df.head(max_rows)
    rows, cols = data.shape
    if rows == 0 or cols == 0:
        return
    total_h = min(height, 0.42 + (rows + 1) * 0.30)
    graphic = slide.shapes.add_table(
        rows + 1, cols, Inches(left), Inches(top), Inches(width), Inches(total_h)
    )
    table = graphic.table
    lens = []
    for j, col in enumerate(data.columns):
        m = sum(2 if ord(ch) > 127 else 1 for ch in str(col))
        for value in data[col]:
            s = _display_value(col, value)
            m = max(m, sum(2 if ord(ch) > 127 else 1 for ch in s))
        lens.append(max(m, 5))
    total_len = sum(lens)
    raw_widths = [max(width * ln / total_len, 0.6) for ln in lens]
    scale = width / sum(raw_widths)
    for j in range(cols):
        table.columns[j].width = Inches(raw_widths[j] * scale)
    table.rows[0].height = Inches(0.38)
    for i in range(1, rows + 1):
        table.rows[i].height = Inches(0.30)
    for j, col in enumerate(data.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (_, row) in enumerate(data.iterrows(), start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = _display_value(data.columns[j], value)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table


def _add_chart(slide, chart_type, title, categories, values,
               left=CHART_LEFT, top=CHART_TOP, width=CHART_W, height=CHART_H, legend=False):
    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    chart_data.add_series("数值", _chart_values(values))
    chart = slide.shapes.add_chart(
        chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    ).chart
    chart.has_legend = legend
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
    return chart


def _metrics_table(df):
    if df is None or df.empty:
        return pd.DataFrame()
    row = df.iloc[0]
    out = pd.DataFrame({"指标": row.index, "数值": row.values})
    # 时间类指标（指标名含 秒/时长）：<60s 按秒、>=60s 按 h:mm:ss 展示
    mask = out["指标"].astype(str).str.contains("秒|时长", regex=True, na=False)
    for idx in out.index[mask]:
        v = out.at[idx, "数值"]
        try:
            out.at[idx, "数值"] = fmt_duration(float(v), with_ms=False)
        except (TypeError, ValueError):
            pass
    return out


def _find_cn_font():
    """查找系统中文字体（macOS 苹方/黑体、Windows 微软雅黑/黑体）。"""
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    return None


def _gantt_page_rows(uph_path):
    """
    从 UPH Excel 的「步骤甘特图」sheet 筛选瓶颈工序的甘特行：
    - CAW（瓶颈机台=焊接机）：取表现居中的滑台（4 个滑台按周期中位排序取中间），
      返回该滑台的 左+右 工位步骤行，标题含滑台号；
    - SA（瓶颈工位）：返回全部工位轨道行（并行对比，标题标瓶颈工位）；
    - 其他制程：返回全部行。
    """
    gantt = _read_sheet(uph_path, "步骤甘特图")
    if gantt.empty or '单元' not in gantt.columns:
        return pd.DataFrame(), ''
    for col in ('开始秒', '结束秒', '时长秒'):
        if col in gantt.columns:
            gantt[col] = pd.to_numeric(gantt[col], errors='coerce')
    ame = _read_sheet(uph_path, "AMESummary")
    if not ame.empty and "瓶颈机台" in ame.columns:
        b_name = str(ame.iloc[0].get("瓶颈机台") or "")
        if b_name and b_name in "焊接机":
            # 收集瓶颈机台的滑台工位单元，按滑台聚合周期（左右结束秒最大值）
            prefix = "焊接机-滑台"
            units = [u for u in gantt['单元'].unique()
                     if str(u).startswith(prefix)]
            slide_rows = {}
            for u in units:
                m = str(u)[len(prefix):]
                if m[:-1].isdigit():
                    n = int(m[:-1])
                    sub = gantt[gantt['单元'] == u]
                    cyc = sub['结束秒'].max() if not sub.empty else 0.0
                    slide_rows.setdefault(n, []).append(cyc)
            if slide_rows:
                # 滑台周期 = 左右工位中较慢者；取周期中位对应的滑台
                order = sorted(slide_rows, key=lambda n: max(slide_rows[n]))
                sel = order[len(order) // 2]
                sel_rows = gantt[gantt['单元'].isin(
                    ["焊接机-滑台%d左" % sel, "焊接机-滑台%d右" % sel])].copy()
                title = f"瓶颈机台：{b_name}（居中滑台：滑台{sel}，左右工位）"
                return sel_rows, title
    if not ame.empty and "瓶颈工位" in ame.columns:
        b_name = str(ame.iloc[0].get("瓶颈工位") or "")
        title = f"瓶颈工位：{b_name}（四工位并行甘特）" if b_name else "工位甘特图"
        return gantt, title
    return gantt, "步骤甘特图（整机）"


def _gantt_png(rows_df, title, out_path):
    """PIL 绘制步骤甘特图 PNG（深色驾驶舱样式，与模板一致）。"""
    if rows_df is None or rows_df.empty:
        return False
    font_path = _find_cn_font()
    def font(size):
        return ImageFont.truetype(font_path, size) if font_path else ImageFont.load_default()
    f_title = font(26)
    f_label = font(17)
    f_axis = font(14)
    bg = (13, 21, 36)          # 0D1524
    grid = (36, 51, 78)        # 24334E
    axis = (143, 163, 196)     # 8FA3C4
    title_c = (231, 241, 255)  # E7F1FF
    label_c = (198, 211, 232)  # C6D3E8
    colors = {
        '循环': (53, 198, 244),    # 青
        '批次': (245, 182, 76),    # 琥珀
        '盘': (46, 134, 193),      # 钢蓝
    }
    n = len(rows_df)
    left_pad, right_pad, top_pad, bottom_pad = 280, 40, 78, 46
    row_h = 30
    max_end = max(float(rows_df['结束秒'].max()), 1.0)
    plot_w = 1180
    px_per_sec = plot_w / max_end
    img_h = top_pad + n * row_h + bottom_pad
    img = Image.new("RGB", (left_pad + plot_w + right_pad, img_h), bg)
    d = ImageDraw.Draw(img)
    d.text((20, 18), title, fill=title_c, font=f_title)
    # 网格与刻度
    step = 1.0
    while max_end / step > 24:
        step *= 2
    for s in range(0, int(math.ceil(max_end / step)) + 1):
        x = left_pad + s * step * px_per_sec
        d.line([(x, top_pad - 14), (x, top_pad + n * row_h)], fill=grid)
        d.text((x - 8, top_pad - 34), str(int(s * step)), fill=axis, font=f_axis)
    d.line([(left_pad, top_pad + n * row_h), (left_pad + plot_w, top_pad + n * row_h)],
           fill=axis, width=2)
    # 步骤横条
    for i, (_, r) in enumerate(rows_df.iterrows()):
        y = top_pad + i * row_h
        label = "%s · %s" % (r['单元'], r['步骤'])
        d.text((8, y + 4), label, fill=label_c, font=f_label)
        start = float(r['开始秒'])
        end = float(r['结束秒'])
        x0 = left_pad + start * px_per_sec
        x1 = left_pad + end * px_per_sec
        if x1 - x0 < 3:
            x1 = x0 + 3
        color = colors.get(str(r['层级']), colors['循环'])
        d.rectangle([x0, y + 2, x1, y + row_h - 4], fill=color, outline=(24, 33, 51))
    img.save(out_path)
    return True


def _report_title(process_name):
    prefix = _TITLE_PREFIX.get(process_name)
    if prefix:
        return f"{prefix}設備一鍵自動分析報告"
    return "MC Log 分析报告"


def _fill_template_cover(prs, process_name):
    """替换封面标题中的制程字母（XX/LM/CAW/FR），其余内容与样式保持不变。

    兼容两种模板形态：
    - 旧模板：标题为独立 run（"LM" + "設備一鍵自動分析報告"）；
    - 新模板（成品视觉稿）：标题为整段单 run（"LM設備一鍵自動分析報告"）。
    """
    if not prs.slides:
        return
    prefix = _TITLE_PREFIX.get(process_name)
    if prefix is None:
        return
    cover = prs.slides[0]
    for shape in cover.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if "設備一鍵自動分析報告" not in tf.text:
            continue
        if tf.paragraphs and tf.paragraphs[0].runs:
            runs = tf.paragraphs[0].runs
            old = "".join(r.text for r in runs)
            head = old.split("設備")[0].replace(" ", "")
            if head in ("XX", "LM", "CAW", "FR", "SA", "ACF"):
                runs[0].text = prefix + old[len(head):]
                for r in runs[1:]:
                    r.text = ""
        return


def _delete_slide(prs, index):
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    sld_id = slides[index]
    r_id = sld_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(r_id)
    sld_id_lst.remove(sld_id)


def _delete_slide_by_title(prs, title):
    sld_id_lst = prs.slides._sldIdLst
    for idx in range(len(sld_id_lst) - 1, 0, -1):
        slide = prs.slides[idx]
        if _slide_has_title(slide, title):
            _delete_slide(prs, idx)
            return


def _prepare_template_pages(prs):
    """保留封面与分节内容页，删除目录页及不属于分节的多余页面。"""
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    for idx in range(len(slides) - 1, 0, -1):
        slide = prs.slides[idx]
        name = slide.slide_layout.name.lower()
        if 'content' not in name and 'agenda' not in name:
            continue
        t = slide.shapes.title
        title = t.text.strip() if t is not None else ''
        if title in SECTION_TITLES:
            continue
        _delete_slide(prs, idx)


def _section_slide(prs, title):
    for slide in prs.slides:
        if _slide_has_title(slide, title):
            return slide
    return None


def _slide_has_title(slide, title):
    """按任意文本框内容匹配分节标题（兼容占位符标题与成品稿文本框标题）。"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == title:
            return True
    return False


def _set_subtitle(slide, subtitle):
    """更新副标题：优先命中标题下方通栏文本框，兼容成品视觉稿模板。"""
    for shape in slide.shapes:
        if shape.is_placeholder or shape.has_chart or shape.has_table or not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if not tf.text.strip():
            continue
        top = (shape.top or 0) / 914400.0
        width = (shape.width or 0) / 914400.0
        if 0.9 <= top <= 2.3 and width >= 4.0:
            _set_text_preserving(tf, subtitle)
            return
    for shape in slide.shapes:
        if shape.is_placeholder or shape.has_chart or shape.has_table or not shape.has_text_frame:
            continue
        _set_text_preserving(shape.text_frame, subtitle)
        return


def _replace_chart(slide, chart_type, chart_title, categories, values, legend=False):
    for shape in list(slide.shapes):
        if not shape.has_chart:
            continue
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        shape._element.getparent().remove(shape._element)
        chart_data = CategoryChartData()
        chart_data.categories = list(categories)
        chart_data.add_series("数值", _chart_values(values))
        chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart
        chart.has_legend = legend
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
        return chart
    return None


def _refresh_chart(slide, chart_type, chart_title, categories, values):
    """原位刷新模板图表数据（replace_data），保留模板图表样式/颜色/标签/图例。

    与 _replace_chart（删除重建）不同，不会丢深色驾驶舱模板的自定义样式，
    也不会在包内残留孤儿图表部件。仅当类型不匹配等极端情况回退到重建。
    """
    for shape in slide.shapes:
        if not shape.has_chart:
            continue
        chart_data = CategoryChartData()
        chart_data.categories = list(categories)
        chart_data.add_series("数值", _chart_values(values))
        try:
            shape.chart.replace_data(chart_data)
            return shape.chart
        except Exception:
            return _replace_chart(slide, chart_type, chart_title, categories, values)
    return None


def _remove_table(slide):
    """删除页面上的表格（瓶颈模式 PPT 只放最终结果，不展示工位明细）。"""
    for shape in list(slide.shapes):
        if shape.has_table:
            shape._element.getparent().remove(shape._element)


def _update_table(slide, df, max_rows=12):
    """原位更新模板表格，严格受方案C 模板模块边界约束：

    - 行数超出容量时截断；
    - 行高按模块高度均分，保证不溢出模块/页脚；
    - 列数变化时按模块宽度重新均分列宽，避免横向溢出。
    """
    data = df.head(max_rows)
    rows, cols = data.shape
    if rows == 0 or cols == 0:
        return
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        t_el = table._tbl
        # 以模板表格图形框为边界（in）
        box_w = (shape.width or Inches(5)).inches
        box_h = (shape.height or Inches(4.75)).inches
        # 列数不足时扩展（模板表可能只有 2 列，而工位明细有多列）
        need_cols = min(cols, 12)
        grid = t_el.find(qn('a:tblGrid'))
        grid_cols = grid.findall(qn('a:gridCol')) if grid is not None else []
        original_cols = len(grid_cols)
        if len(grid_cols) < need_cols:
            for _ in range(need_cols - len(grid_cols)):
                grid.append(copy.deepcopy(grid_cols[-1]))
            for tr in t_el.findall(qn('a:tr')):
                tcs = tr.findall(qn('a:tc'))
                if tcs:
                    for _ in range(need_cols - len(tcs)):
                        tr.append(copy.deepcopy(tcs[-1]))
        elif len(grid_cols) > need_cols:
            for gc in grid_cols[need_cols:]:
                grid.remove(gc)
            for tr in t_el.findall(qn('a:tr')):
                for tc in tr.findall(qn('a:tc'))[need_cols:]:
                    tr.remove(tc)
        # 行数容量：模块高度 / 最小行高 - 表头
        capacity = max(1, int(box_h / _TEMPLATE_DESIGN["table_min_row_in"]) - 1)
        capacity = min(capacity, _TEMPLATE_DESIGN["table_max_data_rows"])
        if rows > capacity:
            data = data.head(capacity)
            rows = capacity
        trs = t_el.findall(qn('a:tr'))
        need = rows + 1
        if len(trs) < need:
            last_tr = trs[-1]
            for _ in range(need - len(trs)):
                t_el.append(copy.deepcopy(last_tr))
        elif len(trs) > need:
            for tr in trs[need:]:
                t_el.remove(tr)
        # 行高均分到模块高度（不溢出）
        row_h_in = box_h / need
        for ri in range(need):
            table.rows[ri].height = Inches(row_h_in)
        # 列数变化时按模块宽度重新均分；不变时保留模板原始列宽
        if need_cols != original_cols and need_cols > 0:
            col_w_in = box_w / need_cols
            for ci in range(need_cols):
                table.columns[ci].width = Inches(col_w_in)
        n_cols = min(cols, need_cols)
        for j in range(n_cols):
            _set_text_preserving(table.cell(0, j).text_frame, str(data.columns[j]))
        for i, (_, row) in enumerate(data.iterrows(), start=1):
            for j in range(n_cols):
                value = _display_value(data.columns[j], row.iloc[j])
                _set_text_preserving(table.cell(i, j).text_frame, value)
        return


def _build_section(prs, title, subtitle, chart_type, chart_title, cats, vals, table_df,
                   table_left=TABLE_LEFT, table_width=TABLE_W, chart_full=False):
    """优先原位更新模板内容页；模板无此页时新建。"""
    slide = _section_slide(prs, title)
    if slide is None:
        slide = _add_slide(prs, title, subtitle)
        if cats is not None and len(cats) > 0:
            _add_chart(slide, chart_type, chart_title, cats, vals,
                       width=(12.1 if chart_full else CHART_W),
                       legend=(chart_type == XL_CHART_TYPE.PIE))
        if not table_df.empty:
            _add_table(slide, table_df, left=table_left, top=CHART_TOP,
                       width=(12.1 if chart_full else table_width), height=CHART_H)
        return
    _set_subtitle(slide, subtitle)
    if cats is not None and len(cats) > 0:
        _refresh_chart(slide, chart_type, chart_title, cats, vals)
    if not table_df.empty:
        _update_table(slide, table_df)


def _update_page_numbers(prs):
    """同步页码：成品稿模板为合成文本框（1/7、2/7…），旧模板为占位符（'1 /' + 总数）。"""
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        page_holder = page_x = None
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            idx = shape.placeholder_format.idx
            if idx == 4:
                page_holder = shape
            elif idx == 3:
                page_x = shape
        if page_holder is not None:
            _set_text_preserving(page_holder.text_frame, f"{i} /")
        if page_x is not None:
            if page_holder is not None:
                _set_text_preserving(page_x.text_frame, str(total))
            else:
                _set_text_preserving(page_x.text_frame, f"{i}/{total}")
        if page_holder is None and page_x is None and not _update_page_textbox(slide, i, total):
            tb = slide.shapes.add_textbox(Inches(11.35), Inches(7.08), Inches(1.75), Inches(0.3))
            tf = tb.text_frame
            tf.text = f"{i}/{total}"
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = GRAY


def _update_page_textbox(slide, index, total):
    """成品视觉稿模板：更新右下角合成页码文本框（1/7、2/7…），返回是否命中。"""
    pattern = re.compile(r"^\d{1,2}/\d{1,2}$")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if pattern.match(txt):
            _set_text_preserving(shape.text_frame, f"{index}/{total}")
            return True
    return False


def _update_section_chips(slide, index, total):
    """同步页头胶囊页码（02 / 07 → 实际页序），删页后保持一致。"""
    pattern = re.compile(r"^\d{2} / \d{2}$")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if pattern.match(txt):
            _set_text_preserving(shape.text_frame, f"{index:02d} / {total:02d}")
            return


def build_ppt_report(output_dir, process_name=None, report_name="Analysis_Report.pptx"):
    """读取 output_dir 下 4 个分析 Excel，基于最终模板生成报告，返回报告路径。"""
    out_path = os.path.join(output_dir, report_name)
    template = next(
        (p for p in (find_external_resource(name) for name in _TEMPLATE_CANDIDATES) if p),
        None,
    )
    prs = Presentation(template) if template else Presentation()
    if template:
        _fill_template_cover(prs, process_name)
        _prepare_template_pages(prs)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = _add_slide(prs, _report_title(process_name), "UPH / EFF / 报警 / 机台状态")
        box = slide.shapes.add_textbox(Inches(0.4), Inches(3.2), Inches(9.2), Inches(0.5))
        box.text_frame.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        box.text_frame.paragraphs[0].font.size = Pt(14)
        box.text_frame.paragraphs[0].font.color.rgb = GRAY

    to_remove = []

    # ---------- UPH 分析 ----------
    ame = _read_sheet(os.path.join(output_dir, "UPH_Analysis.xlsx"), "AMESummary")
    if ame.empty:
        to_remove.append("UPH 分析")
    else:
        uph_summary = _read_sheet(os.path.join(output_dir, "UPH_Analysis.xlsx"), "Summary")
        machine_mode = not ame.empty and "瓶颈机台" in ame.columns
        if machine_mode:
            # CAW 双机台：瓶颈机台 + 单颗CT + 整机 UPH，各机台 CT 用图展示
            b_name = str(ame.iloc[0].get("瓶颈机台") or "")
            b_ct = str(ame.iloc[0].get("瓶颈CT(秒)") or "")
            uph = str(ame.iloc[0].get("UPH(个/小时)") or "")
            subtitle = f"瓶颈机台：{b_name}（单颗CT {b_ct}s）→ UPH ≈ {uph}/h"
            slide = _section_slide(prs, "UPH 分析")
            if slide is None:
                slide = _add_slide(prs, "UPH 分析", subtitle)
            else:
                _set_subtitle(slide, subtitle)
            machine_table = uph_summary.copy()
            if "单颗CT(秒)" in machine_table.columns:
                machine_table["单颗CT(秒)"] = machine_table["单颗CT(秒)"].fillna("")
            _update_table(slide, machine_table)
            _refresh_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, "各机台单颗CT（秒）",
                           uph_summary["机台"], uph_summary["单颗CT(秒)"])
        elif not uph_summary.empty and "工位" in uph_summary.columns:
            # 多工位机：PPT 只放最终结果——瓶颈工位与 UPH，各工位周期用图展示
            b_name = str(ame.iloc[0].get("瓶颈工位") or "")
            b_time = str(ame.iloc[0].get("瓶颈周期(秒)") or "")
            pure = str(ame.iloc[0].get("Pure UPH(个/小时)") or "")
            subtitle = f"瓶颈工位：{b_name}（{b_time}s/排）→ UPH ≈ {pure}/h"
            slide = _section_slide(prs, "UPH 分析")
            if slide is None:
                slide = _add_slide(prs, "UPH 分析", subtitle)
            else:
                _set_subtitle(slide, subtitle)
            station_table = uph_summary.copy()
            if "瓶颈" in station_table.columns:
                station_table["瓶颈"] = station_table["瓶颈"].fillna("")
            _update_table(slide, station_table)
            _refresh_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, "各工位每排周期（秒）",
                           station_table["工位"], station_table["每排周期(秒)"])
        else:
            labels = ["Pure UPH(个/小时)", "Derated UPH M1(个/小时)", "Derated UPH M2(个/小时)"]
            vals = [_to_float(ame.iloc[0].get(c)) for c in labels]
            _build_section(prs, "UPH 分析", "CoreTech AME：Pure UPH / Derated UPH M1 / M2",
                           XL_CHART_TYPE.COLUMN_CLUSTERED, "UPH 对比", labels, vals, _metrics_table(ame))

    # ---------- 瓶颈工序甘特图 ----------
    gantt_rows, gantt_title = _gantt_page_rows(os.path.join(output_dir, "UPH_Analysis.xlsx"))
    if not gantt_rows.empty:
        png_path = os.path.join(output_dir, "_gantt_bottleneck.png")
        try:
            if _gantt_png(gantt_rows, gantt_title, png_path):
                slide = _add_dark_gantt_slide(prs, len(prs.slides) + 1, len(prs.slides) + 1)
                st = slide.shapes.add_textbox(Inches(0.53), Inches(1.22), Inches(12.27), Inches(0.4))
                st.text_frame.text = gantt_title
                st.text_frame.paragraphs[0].font.size = Pt(18)
                st.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x8F, 0xA3, 0xC4)
                with Image.open(png_path) as img:
                    # Windows 下 Image 对象未关闭会持有文件句柄，
                    # 导致下方 os.remove 报 WinError 32，必须用 with 及时释放。
                    ratio = img.width / img.height
                w = 12.13
                h = w / ratio
                if h > 4.55:
                    h = 4.55
                    w = h * ratio
                slide.shapes.add_picture(png_path,
                                         Inches(0.4 + (12.53 - w) / 2),
                                         Inches(1.78 + (4.75 - h) / 2),
                                         width=Inches(w), height=Inches(h))
        finally:
            try:
                if os.path.exists(png_path):
                    os.remove(png_path)
            except OSError:
                # 个别环境（杀毒软件/索引服务）可能瞬时占用文件，
                # 删除失败不应让整个 PPT 报告生成中断。
                pass

    # ---------- EFF 分析 ----------
    eff = _read_sheet(os.path.join(output_dir, "EFF_Analysis.xlsx"), "Summary")
    if eff.empty:
        to_remove += ["EFF 分析", "停机 Pareto"]
    else:
        row = eff.iloc[0]
        cats = ["运行RUN", "待机IDLE", "停机DOWN"]
        vals = [_to_float(row.get(c)) for c in ("运行时间RUN(秒)", "待机时间IDLE(秒)", "停机时间DOWN(秒)")]
        _build_section(prs, "EFF 分析", "EFF = 操作时间(运行+待机) / 计划生产时间",
                       XL_CHART_TYPE.PIE, "时间构成（秒）", cats, vals, _metrics_table(eff))
        pareto = _read_sheet(os.path.join(output_dir, "EFF_Analysis.xlsx"), "DOWN_Pareto")
        if pareto.empty:
            to_remove.append("停机 Pareto")
        else:
            top = pareto.head(7)  # 只放前 7 项，避免超过 PPT 页面
            if '原因名称' in top.columns and top['原因名称'].notna().any():
                # 有原因清单：以中文原因名称展示
                labels = top['原因名称'].fillna('')
                table_df = top[['原因名称', '次数', '总时长(秒)', '占比(%)']].copy()
                subtitle = "按停机原因统计时长 Top10"
            else:
                labels = top['ReasonID']
                table_df = top[['ReasonID', '次数', '总时长(秒)', '占比(%)']]
                subtitle = "按 ReasonID 统计停机时长 Top10"
            _build_section(prs, "停机 Pareto", subtitle,
                           XL_CHART_TYPE.BAR_CLUSTERED, "停机时长（秒）",
                           labels, top["总时长(秒)"], table_df)

    # ---------- 报警分析 ----------
    by_kw = _read_sheet(os.path.join(output_dir, "Alarm_Analysis.xlsx"), "ByKeyword")
    alarm_summary = _read_sheet(os.path.join(output_dir, "Alarm_Analysis.xlsx"), "Summary")
    if by_kw.empty and alarm_summary.empty:
        to_remove.append("报警分析")
    else:
        if not by_kw.empty:
            top = by_kw.head(10)
            _build_section(prs, "报警分析", "按关键词与模块统计报警",
                           XL_CHART_TYPE.COLUMN_CLUSTERED, "报警次数 Top10",
                           top["命中关键词"], top["报警次数"], top)
        else:
            _build_section(prs, "报警分析", "按关键词与模块统计报警",
                           XL_CHART_TYPE.COLUMN_CLUSTERED, "报警次数 Top10",
                           [], [], alarm_summary)

    # ---------- 机台状态 ----------
    status = _read_sheet(os.path.join(output_dir, "Status_Analysis.xlsx"), "Summary")
    if status.empty:
        to_remove += ["机台状态汇总", "机台状态趋势"]
    else:
        cats = list(status["状态"])
        vals = [_to_float(v) for v in status["总时长(秒)"]]
        _build_section(prs, "机台状态汇总", "RUN / IDLE / DOWN / WARN 时长与占比",
                       XL_CHART_TYPE.PIE, "状态时长（秒）", cats, vals, status)
        hourly = _read_sheet(os.path.join(output_dir, "Status_Analysis.xlsx"), "Hourly")
        if hourly.empty or "EFF(%)" not in hourly.columns:
            to_remove.append("机台状态趋势")
        else:
            _build_section(prs, "机台状态趋势", "每小时 EFF 趋势",
                           XL_CHART_TYPE.LINE_MARKERS, "每小时 EFF(%)",
                           hourly["小时"], hourly["EFF(%)"], pd.DataFrame(), chart_full=True)

    for title in to_remove:
        _delete_slide_by_title(prs, title)

    _update_page_numbers(prs)
    for i, slide in enumerate(prs.slides, start=1):
        _update_section_chips(slide, i, len(prs.slides))
    prs.save(out_path)
    return out_path
