#!/usr/bin/env python3
"""MC Log Analysis 使用说明 deck 构建脚本（基于 PPT模板.pptx，master-first）。

继承模板的 master/layout（版权、logo、页码/页脚、封面/目录/内容页体系），
按模板配色（强调 EF5149、深蓝 083C63）排版，输出可编辑 pptx。
"""

from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SPEC_PATH = ROOT / "build" / "generated" / "slide_specs.yaml"
TEMPLATE = ROOT.parents[1] / "PPT模板.pptx"
OUT_PATH = ROOT / "build" / "pptx" / "usage_deck.pptx"
GANTT_PNG = ROOT / "assets" / "images" / "generated" / "gantt_example_lm.png"

ACCENT = RGBColor(0xEF, 0x51, 0x49)   # 模板强调红橙（封面标题）
NAVY = RGBColor(0x08, 0x3C, 0x63)     # 模板深蓝（页码/次级标题）
GRAY = RGBColor(0x59, 0x59, 0x59)
GRAY_LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x26, 0x26, 0x26)
LINE_GRAY = RGBColor(0xD9, 0xD9, 0xD9)

EA_FONT = "宋体"
LATIN_FONT = "Times New Roman"
TOTAL_SLIDES = 13


def _set_run(run, size, bold=False, color=BLACK, latin=None, ea=None):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = latin or LATIN_FONT
    rPr = run._r.get_or_add_rPr()
    ea_el = rPr.find(qn("a:ea"))
    if ea_el is None:
        ea_el = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea_el)
    ea_el.set("typeface", ea or EA_FONT)


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             space_after=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        _set_run(r, size, bold, color)
    return box


def add_rect_text(slide, x, y, w, h, text, fill=WHITE, line_color=None, radius=None,
                  size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, margin=0.14):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        _set_run(r, size, bold, color)
    return shp


def add_line(slide, x, y, w, color=NAVY, weight=1.5):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_arrow(slide, x, y, w, color=NAVY, weight=2.0):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 Inches(x + w - 0.06), Inches(y - 0.07),
                                 Inches(0.14), Inches(0.14))
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    tri.shadow.inherit = False


def _delete_slide(prs, index):
    sld_id_lst = prs.slides._sldIdLst
    sld_id = list(sld_id_lst)[index]
    r_id = sld_id.get(qn("r:id"))
    prs.part.drop_rel(r_id)
    sld_id_lst.remove(sld_id)


def new_slide(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_title(slide, eyebrow, title):
    """内容页标题占位：章节小字 + 主标题（深蓝）。"""
    ph = slide.shapes.title
    ph.left = Inches(0.64)
    ph.top = Inches(0.25)
    ph.width = Inches(12.0)
    ph.height = Inches(1.15)
    tf = ph.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = eyebrow
    _set_run(r0, 10.5, True, ACCENT)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(2)
    r1 = p1.add_run()
    r1.text = title
    _set_run(r1, 24, True, NAVY)
    return ph


def content_slide(prs, eyebrow, title):
    slide = new_slide(prs, 6)  # 1_content page
    set_title(slide, eyebrow, title)
    return slide


def build_cover(prs):
    slide = new_slide(prs, 0)  # 1_cover
    add_text(slide, 0.88, 2.5, 12.3, 1.1, "MC Log Analysis",
             size=54, bold=True, color=ACCENT)
    add_text(slide, 0.9, 3.75, 12.0, 0.5, "机台日志分析工具 · 使用说明",
             size=16, bold=True, color=NAVY)
    add_line(slide, 0.92, 4.5, 3.4, color=NAVY, weight=2.0)
    add_text(slide, 0.92, 4.75, 12.0, 0.4,
             "功能与亮点  ·  操作教学  ·  制程与口径  ·  输出与注意", size=14, color=GRAY)
    add_text(slide, 0.92, 6.65, 8.0, 0.3, "v1.0  ·  内部使用", size=9, color=GRAY)


def build_agenda(prs):
    slide = new_slide(prs, 1)  # 1_agenda
    ph = slide.shapes.title
    ph.text_frame.text = "目录"
    body = slide.placeholders[10]
    body.text_frame.clear()
    items = [
        "第1章  功能与亮点：工具能做什么、亮点在哪",
        "第2章  操作教学：5 步完成一次分析、界面导览",
        "第3章  制程与口径：制程模板、UPH/EFF 计算逻辑、步骤甘特图",
        "第4章  输出与注意：Excel/PPT 输出物、使用注意事项",
    ]
    for i, it in enumerate(items):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.line_spacing = 1.4
        r = p.add_run()
        r.text = it
        _set_run(r, 14, False, NAVY if i == 0 else BLACK)


def build_overview(prs):
    slide = content_slide(prs, "第1章 功能与亮点", "一个工具，完成机台日志的六类分析与报告")
    add_text(slide, 0.7, 1.62, 12.0, 0.5,
             "把机台日志的“读、算、查、报”一次做完：选择日志文件夹与制程，一键输出 Excel 与 PPT 报告。",
             size=14, color=GRAY)
    nums = [
        ("6", "类分析功能", "合并拆分 · UPH · EFF · 报警 · 状态 · 一键"),
        ("5", "种制程模板", "LM / CAW / FR / SA / ACF，参数自动预填"),
        ("4+1", "类输出", "UPH / EFF / 报警 / 状态 Excel + 自动 PPT 报告"),
    ]
    for i, (num, label, desc) in enumerate(nums):
        x = 0.7 + i * 4.1
        add_text(slide, x, 2.7, 3.8, 1.0, num, size=40, bold=True, color=ACCENT)
        add_text(slide, x + 0.05, 3.85, 3.7, 0.4, label, size=16, bold=True, color=NAVY)
        add_text(slide, x + 0.05, 4.35, 3.75, 1.0, desc, size=12, color=GRAY, line_spacing=1.3)
    add_rect_text(slide, 0.7, 5.9, 12.0, 0.8,
                  "支持“一键分析”：一次运行完成全部功能，结果在界面显示，并自动导出 4 个分析 Excel 与 PPT 报告。",
                  fill=RGBColor(0xFD, 0xEF, 0xEE), size=14, bold=True, color=ACCENT)


def build_grid_page(prs, eyebrow, title, cards, cols=3, fill=NAVY,
                    title_color=WHITE, body_color=RGBColor(0xDD, 0xE6, 0xEF)):
    slide = content_slide(prs, eyebrow, title)
    cw = 3.9 if cols == 3 else 5.95
    ch = 2.25
    gap = 0.14
    for i, (t, body) in enumerate(cards):
        row, col = divmod(i, cols)
        x = 0.7 + col * (cw + gap)
        y = 1.75 + row * (ch + 0.2)
        shp = add_rect_text(slide, x, y, cw, ch, t, fill=fill, size=14, bold=True,
                            color=title_color, anchor=MSO_ANCHOR.TOP, margin=0.14)
        p = shp.text_frame.add_paragraph()
        p.line_spacing = 1.25
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = body
        _set_run(r, 12, False, body_color)
    return slide


def build_flow(prs):
    slide = content_slide(prs, "第2章 操作教学", "5 步完成一次机台日志分析")
    steps = [
        ("1 选择源文件夹", "机台日志所在目录，支持子文件夹"),
        ("2 选择输出文件夹", "Excel / PPT 报告输出位置"),
        ("3 选择制程模板", "自动预填 UPH 触发词、报警关键词等"),
        ("4 选择功能", "合并 / UPH / EFF / 报警 / 状态 / 一键"),
        ("5 自动报告", "界面显示结果，并导出 Excel / PPT"),
    ]
    bw, bh, gap = 2.16, 1.1, 0.38
    y = 2.35
    for i, (t, desc) in enumerate(steps):
        x = 0.7 + i * (bw + gap)
        add_rect_text(slide, x, y, bw, bh, t, fill=NAVY, radius=0.08,
                      size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, x - 0.1, y + bh + 0.16, bw + 0.2, 1.0, desc,
                 size=12, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.3)
        if i < 4:
            add_arrow(slide, x + bw - 0.02, y + bh / 2, gap + 0.06)
    add_text(slide, 0.7, 5.7, 12.0, 0.4,
             "提示：日志文件筛选框会自动填充与制程匹配的内容，有特殊需求也可手动输入。",
             size=12, color=GRAY)


def build_ui_map(prs):
    slide = content_slide(prs, "第2章 操作教学", "界面一看就懂：左边选功能，右边筛日志，下方看结果")
    add_text(slide, 0.7, 1.6, 12.0, 0.4, "界面示意（非截图）：三块区域各司其职。", size=12, color=GRAY)
    add_rect_text(slide, 0.7, 2.15, 2.95, 4.2, "功能选择",
                  fill=RGBColor(0xFD, 0xEF, 0xEE), size=14, bold=True, color=ACCENT,
                  anchor=MSO_ANCHOR.TOP, margin=0.14)
    for j, it in enumerate(["文档合并与内容拆分", "UPH 分析", "EFF 分析", "报警分析", "机台状态分析", "一键分析"]):
        add_text(slide, 0.92, 2.65 + j * 0.58, 2.6, 0.4, "· " + it, size=12, color=GRAY)
    add_rect_text(slide, 3.85, 2.15, 8.85, 2.45, "制程模板 · 日志文件筛选 · 解析参数",
                  fill=GRAY_LIGHT, size=14, bold=True, color=NAVY, anchor=MSO_ANCHOR.TOP, margin=0.14)
    add_text(slide, 4.08, 2.72, 8.4, 0.4, "制程下拉：LM / CAW / FR / SA / ACF / 通用 / 自定义", size=12)
    add_text(slide, 4.08, 3.22, 8.4, 0.4, "日志文件筛选：自动填充与制程匹配的内容，可手动修改", size=12)
    add_text(slide, 4.08, 3.72, 8.4, 0.4, "解析参数：关键词、分隔符（默认左对齐）", size=12)
    add_rect_text(slide, 3.85, 4.85, 8.85, 1.5, "解析结果",
                  fill=WHITE, line_color=NAVY, size=14, bold=True, color=NAVY, anchor=MSO_ANCHOR.TOP, margin=0.14)
    add_text(slide, 4.08, 5.42, 8.4, 0.8,
             "分析结果表格与进度显示；顶部“使用说明”按钮可查看各制程说明与计算逻辑，右下角为版权信息。",
             size=12, color=GRAY, line_spacing=1.4)


def build_process_table(prs):
    slide = content_slide(prs, "第3章 制程与口径", "五种制程模板，参数与算法已按机台配好")
    headers = ["制程", "周期完成事件", "结构", "UPH 口径要点"]
    rows = [
        ["LM 激光打标", "MarkEnd1", "单机台；每盘 4 批 × 6 颗", "换盘时间按 24 颗平摊"],
        ["CAW 组装", "放熟料完成 / 去交换位", "PLC1 焊接机 + PLC2 上料机（左右并行）", "双机台瓶颈，取 CT 长者"],
        ["FR 点胶", "轴点胶完成", "左轴 / 右轴并行", "Pure UPH ×0.5（单轴）；整机 M1/M2 = 左 + 右"],
        ["SA 四工位", "各工位完成事件", "点胶 / 贴附 / 热压 / 检测并行", "自动判定瓶颈工位"],
        ["ACF 三機", "更新Carrier盘 / Cavity cnt:1 / UnloadDuts Finish", "上料機 / 主機 / 下料機 分文件夹", "每盘颗数按 Tray ID 动态统计"],
    ]
    tbl = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.7), Inches(1.75),
                                 Inches(12.0), Inches(3.0)).table
    widths = [2.0, 3.2, 3.6, 3.2]
    for j, w in enumerate(widths):
        tbl.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _set_run(r, 10.5, True, WHITE)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        row_fill = GRAY_LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 1.0
                for r in p.runs:
                    _set_run(r, 10.5, False, BLACK)
            c.fill.solid()
            c.fill.fore_color.rgb = row_fill
    add_text(slide, 0.7, 5.05, 12.0, 0.4,
             "另有“通用（手动配置）”与“自定义”模板，不套用固定逻辑。", size=12, color=GRAY)


def build_formulas(prs):
    slide = content_slide(prs, "第3章 制程与口径", "UPH 口径：实际、Pure、Derated M1/M2 与有效 UPH")
    items = [
        ("实际 UPH = 产出数 ÷ 统计时长 × 3600", "按真实产出与统计时长计算。"),
        ("Pure UPH = 3600 × 每周期产出 ÷ 理想周期CT", "未填理想 CT 时取正常周期平均。"),
        ("Derated UPH M2 = 3600 × 每周期产出 ÷ 有效平均周期", "剔除 <0.9×理想CT、>1.1×最大理论CT 的离群点；多模组整机 = 各模组 M2 之和。"),
        ("Derated UPH M1 = EM 投入数 ÷ RUN 时长", "多模组整机 = 各模组 M1 之和。"),
        ("有效 UPH = 3600 × 每周期产出 ÷（基础周期 + 每颗换盘开销）", "每颗换盘开销 = 单次换盘时间 ÷ 每盘颗数。"),
    ]
    y = 1.7
    for formula, note in items:
        add_text(slide, 0.7, y, 12.0, 0.42, formula, size=14, bold=True, color=NAVY)
        add_text(slide, 0.7, y + 0.42, 12.0, 0.36, note, size=12, color=GRAY)
        y += 1.0
        add_line(slide, 0.7, y - 0.12, 12.0, color=LINE_GRAY, weight=0.75)
    add_text(slide, 0.7, y + 0.05, 9.0, 0.5,
             "周期分类：间隔 > 计划性停机阈值 → 计划性停机；> 正常周期阈值 → 异常周期；其余为正常周期。",
             size=12, color=GRAY)


def build_gantt(prs):
    slide = content_slide(prs, "第3章 制程与口径", "步骤拆到运动节拍，甘特图一眼看出先后与并行")
    add_text(slide, 0.7, 1.6, 6.0, 0.5, "每个区块的步骤时间总和 ≈ CT（周期中位）。",
             size=14, bold=True, color=NAVY)
    add_text(slide, 0.7, 2.15, 6.1, 2.3,
             "步骤按轴 / 平台运动节拍拆分，例如 LM 单颗循环：\n\n"
             "打标前间隔 → Z 轴焦距定位 → 激光打标（振镜扫描）\n\n"
             "批次级：读码 GetSN、CCD 轴移动定位（每批 6 颗）\n\n"
             "盘级：平台下料、平台上料（每盘 24 颗）",
             size=12, color=GRAY, line_spacing=1.35)
    add_text(slide, 0.7, 5.0, 6.1, 1.5,
             "异常判定 = 中位时长 + 超时秒数；超时 / 报警 / 停机行在日志 Excel 中自动标红。",
             size=12, color=GRAY, line_spacing=1.35)
    # 右侧 BITMAP 占位放甘特示例图
    if GANTT_PNG.exists():
        for ph in slide.placeholders:
            try:
                if "BITMAP" in str(ph.placeholder_format.type):
                    ph.insert_picture(str(GANTT_PNG))
                    add_text(slide, 8.0, 6.45, 4.8, 0.35,
                             "甘特示例：来自工具对 LM 示例日志的分析输出（示意）", size=9, color=GRAY)
                    break
            except Exception:
                continue


def build_eff_page(prs):
    cards = [
        ("EFF 分析", "EFF = 操作时间（运行 + 待机）÷ 计划生产时间；停机按 ReasonID 拆 pDT / uDT（EReason 清单优先）。"),
        ("机台状态分析", "优先读 status:RUN/IDLE/DOWN 行；无 status 时按“活动 + 停机关键词”推导时间线，输出各状态时长 / 占比 / 小时分布。"),
        ("报警分析", "按关键词命中计数，按机台 / 模组汇总；EReason 清单映射中文原因名称。"),
    ]
    return build_grid_page(prs, "第3章 制程与口径", "EFF、状态与报警：三张表讲清设备效率与异常",
                           cards, cols=3, fill=NAVY)


def build_output_table(prs):
    slide = content_slide(prs, "第4章 输出与注意", "一次分析，输出可编辑的 Excel 与 PPT 报告")
    headers = ["文件", "内容"]
    rows = [
        ["UPH_Analysis.xlsx", "Summary / AMESummary / CycleDetail / EMProduction / 步骤分析 / 步骤甘特图"],
        ["EFF_Analysis.xlsx", "EFF 汇总与停机 Pareto（含 pDT / uDT）"],
        ["Alarm_Analysis.xlsx", "报警汇总、关键词分布、明细（EReason 中文名）"],
        ["Status_Analysis.xlsx", "状态汇总与小时分布"],
        ["LogAnalysis.xlsx", "合并日志（异常行标红）；大文件按日志逐个导出"],
        ["Analysis_Report.pptx", "自动报告：UPH / EFF / 停机 Pareto / 报警 / 状态 + 瓶颈工序甘特图"],
    ]
    tbl = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.7), Inches(1.75),
                                 Inches(12.0), Inches(3.2)).table
    tbl.columns[0].width = Inches(3.4)
    tbl.columns[1].width = Inches(8.6)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _set_run(r, 10.5, True, WHITE)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        row_fill = GRAY_LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = 1.0
                for r in p.runs:
                    _set_run(r, 10.5, False, BLACK)
            c.fill.solid()
            c.fill.fore_color.rgb = row_fill
    add_text(slide, 0.7, 5.25, 12.0, 0.4,
             "时间格式约定：低于 60 秒按秒显示，超过 60 秒按 h:mm:ss；PPT 汇总时间显示到秒。",
             size=12, color=GRAY)


def build_tips(prs):
    cards = [
        ("大日志保护", "原始日志超过约 100 万行时，自动放弃合并、按日志文件逐个导出 Excel，避免卡死。"),
        ("EReasonlist", "根目录保留 EReasonlist 文件夹，新制程按文件名放入自己的清单，用于状态推导与报警映射。"),
        ("PPT 模板", "Analysis_Report.pptx 放到程序同目录即可正常生成报告；模板含内部数据，仅限内部传送。"),
        ("灵活覆盖", "日志文件筛选框自动填充与制程匹配的内容，用户有特殊需求可手动输入。"),
    ]
    return build_grid_page(prs, "第4章 输出与注意", "使用前注意这几点，避免踩坑",
                           cards, cols=2, fill=NAVY)


def build_closing(prs):
    slide = new_slide(prs, 0)  # 1_cover，与封面呼应
    add_text(slide, 0.9, 2.6, 12.0, 1.0, "现在就可以开始第一次分析",
             size=40, bold=True, color=ACCENT)
    add_text(slide, 0.92, 3.85, 12.0, 0.5, "选日志文件夹 → 选制程模板 → 点“自动报告”",
             size=16, bold=True, color=NAVY)
    add_line(slide, 0.92, 4.6, 3.4, color=NAVY, weight=2.0)
    add_text(slide, 0.92, 4.9, 12.0, 0.4,
             "功能与亮点  ·  操作教学  ·  制程与口径  ·  输出与注意", size=14, color=GRAY)
    add_text(slide, 0.92, 6.65, 8.0, 0.3, "v1.0  ·  内部使用", size=9, color=GRAY)


def fix_docprops_slide_count(path):
    import re
    import shutil
    import zipfile
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "docProps/app.xml":
                s = data.decode("utf-8")
                s = re.sub(r"<Slides>\s*\d*\s*</Slides>", "<Slides>%d</Slides>" % TOTAL_SLIDES, s)
                if "<Slides>" not in s:
                    s = s.replace("</Properties>", "<Slides>%d</Slides></Properties>" % TOTAL_SLIDES)
                data = s.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(str(tmp), str(path))


def main():
    spec = yaml.safe_load(SPEC_PATH.read_text(encoding="utf-8"))
    prs = Presentation(str(TEMPLATE))
    # 删除模板示例页（封面/目录），用 layout 重建
    while len(prs.slides._sldIdLst) > 0:
        _delete_slide(prs, 0)

    build_cover(prs)
    build_agenda(prs)
    build_grid_page(prs, "第1章 功能与亮点", "六大功能，覆盖日志分析的完整链路", [
        ("文档合并与内容拆分", "按关键词 / 分隔符提取指定行，支持按 PLC 分组导出，异常行自动标红。"),
        ("UPH 分析", "按 CoreTech AME 口径输出实际 / Pure / Derated M1/M2 / 有效 UPH，自动判定瓶颈。"),
        ("EFF 分析", "操作时间 ÷ 计划生产时间；停机按 EReason 清单区分 pDT / uDT。"),
        ("报警分析", "按关键词统计报警，EReason 中文名映射。"),
        ("机台状态分析", "status 行或活动 / 停机关键词推导 RUN / IDLE / DOWN 时间线。"),
        ("一键分析", "一次运行完成全部功能，导出 4 个 Excel 并自动生成 PPT 报告。"),
    ], fill=NAVY)
    build_grid_page(prs, "第1章 功能与亮点",
                    "亮点：自动模板、瓶颈判定、换盘分摊、步骤甘特、异常标红、自动报告", [
        ("制程模板一键预填", "选 LM / CAW / FR / SA / ACF，自动带入触发词、报警关键词与计算逻辑。"),
        ("瓶颈自动判定", "SA 四工位、CAW 双机台自动取最慢环节，UPH 口径与客户定义一致。"),
        ("换盘时间平摊", "整盘结束后的下料-上料时间按每盘颗数摊入单颗 CT。"),
        ("步骤深度分析 + 甘特图", "按轴 / 平台运动节拍拆步骤，Excel 甘特图直接看出先后与并行。"),
        ("异常日志自动标红", "报警 / 停机 / 步骤超时行标红，剔除产品码中的 NG 误报。"),
        ("自动 PPT 报告", "汇总 + 图表 + 瓶颈工序甘特图，时间格式统一、可编辑。"),
    ], fill=ACCENT, body_color=RGBColor(0xFF, 0xEF, 0xED))
    build_flow(prs)
    build_ui_map(prs)
    build_process_table(prs)
    build_formulas(prs)
    build_gantt(prs)
    build_eff_page(prs)
    build_output_table(prs)
    build_tips(prs)
    build_closing(prs)

    # 模板 SLIDE_NUMBER 占位符在 add_slide 时未实例化（PowerPoint 渲染无页码），
    # 统一在每页右下角补页码（模板页码色 083C63）。
    for i, slide in enumerate(prs.slides, 1):
        add_text(slide, 11.35, 7.05, 0.72, 0.3, "%d / %d" % (i, TOTAL_SLIDES),
                 size=9, color=NAVY, align=PP_ALIGN.RIGHT)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    fix_docprops_slide_count(OUT_PATH)
    print("saved:", OUT_PATH)


if __name__ == "__main__":
    main()
