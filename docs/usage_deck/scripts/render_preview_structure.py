#!/usr/bin/env python3
"""结构近似预览渲染器（PIL）。

环境限制：本机 PowerPoint 自动化导出 PDF 失败（macOS 自动化权限 -9074），
且无 LibreOffice/pdftoppm，因此用 PIL 按 pptx 的 shape 布局做近似渲染，
用于人工 visual review。渲染结果为“结构近似预览”，非 PowerPoint 高保真成图。
"""

import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "build" / "pptx" / "usage_deck.pptx"
OUT_DIR = ROOT / "build" / "rendered" / "structure_preview"
EMU = 914400
SCALE = 96 / 72  # pt -> px (96 dpi)
TRUNC_WARN = []


def _font(size_pt):
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "C:/Windows/Fonts/msyh.ttc",
    ]
    for p in candidates:
        if Path(p).exists():
            try:
                return ImageFont.truetype(p, max(9, int(size_pt * SCALE)))
            except Exception:
                continue
    return ImageFont.load_default()


def _draw_text(draw, x, y, w, h, text, size_pt, bold, color, align="left"):
    if not text:
        return
    font = _font(size_pt)
    lines = str(text).split("\n")
    line_h = int(size_pt * SCALE * 1.35)
    cur = y + 3
    clipped = False
    for ln in lines:
        if len(lines) > 8:
            clipped = True
            break
        if cur > y + h:
            clipped = True
            break
        if align == "center":
            tw = draw.textlength(ln, font=font)
            tx = x + max(0, (w - tw) / 2)
        elif align == "right":
            tw = draw.textlength(ln, font=font)
            tx = x + max(0, w - tw)
        else:
            tx = x + 4
        draw.text((tx, cur), ln, font=font, fill=color)
        cur += line_h
    if clipped:
        TRUNC_WARN.append(str(text)[:40])


def _render_slide(prs, idx):
    slide = prs.slides[idx]
    W = int(prs.slide_width / EMU * 96)
    H = int(prs.slide_height / EMU * 96)
    img = Image.new("RGB", (W, H), (255, 255, 255))
    d = ImageDraw.Draw(img)
    for shp in slide.shapes:
        x = int(shp.left / EMU * 96)
        y = int(shp.top / EMU * 96)
        w = int(shp.width / EMU * 96)
        h = int(shp.height / EMU * 96)
        st = shp.shape_type
        text = shp.text_frame.text if shp.has_text_frame else ""
        # 背景/填充色
        fill = None
        try:
            if shp.fill.type is not None and str(shp.fill.type) != "MSO_FILL_TYPE.BACKGROUND (5)":
                fill = shp.fill.fore_color.rgb
                fill = (fill[0], fill[1], fill[2])
        except Exception:
            fill = None
        if st == MSO_SHAPE_TYPE.LINE or st == MSO_SHAPE_TYPE.FREEFORM:
            try:
                lc = shp.line.color.rgb
                d.line([(x, y + h // 2), (x + w, y + h // 2)],
                       fill=(lc[0], lc[1], lc[2]), width=2)
            except Exception:
                d.line([(x, y + h // 2), (x + w, y + h // 2)], fill=(60, 60, 60), width=2)
            continue
        if st == MSO_SHAPE_TYPE.PICTURE:
            d.rectangle([x, y, x + w, y + h], outline=(120, 120, 120), width=1)
            d.text((x + 6, y + 6), "【图片】", font=_font(10), fill=(120, 120, 120))
            try:
                from PIL import Image as PILImage
                ph = PILImage.open(shp.image.blob)
                ph.thumbnail((w - 4, h - 4))
                img.paste(ph, (x + 2, y + 2))
            except Exception:
                pass
            continue
        if st == MSO_SHAPE_TYPE.TABLE:
            tbl = shp.table
            nrows, ncols = len(tbl.rows), len(tbl.columns)
            rh = h / nrows
            cw = w / ncols
            for r in range(nrows):
                for c in range(ncols):
                    cx = x + c * cw
                    cy = y + r * rh
                    d.rectangle([cx, cy, cx + cw, cy + rh], outline=(150, 150, 150), width=1)
                    if r == 0:
                        d.rectangle([cx, cy, cx + cw, cy + rh], fill=(222, 235, 247))
                    _draw_text(d, cx, cy, cw, rh, tbl.cell(r, c).text, 10.5, r == 0,
                               (31, 78, 121) if r == 0 else (0, 0, 0))
            continue
        # 形状 / 文本框
        is_round = False
        try:
            is_round = shp.auto_shape_type is not None and "ROUNDED" in str(shp.auto_shape_type)
        except Exception:
            is_round = False
        is_tri = False
        try:
            is_tri = shp.auto_shape_type is not None and "TRIANGLE" in str(shp.auto_shape_type)
        except Exception:
            is_tri = False
        if is_round:
            d.rounded_rectangle([x, y, x + w, y + h], radius=min(12, int(h * 0.12)),
                                fill=fill or (255, 255, 255), outline=(160, 160, 160))
        elif is_tri:
            d.polygon([(x + w, y), (x, y), (x + w, y + h)], fill=fill or (31, 78, 121))
            continue
        else:
            d.rectangle([x, y, x + w, y + h], fill=fill or (255, 255, 255),
                        outline=(160, 160, 160) if fill is None else None)
        # 文本颜色：浅底用黑，深底（深蓝）用白
        tcolor = (255, 255, 255) if fill and sum(fill[:3]) < 350 else (0, 0, 0)
        size_pt = 12
        try:
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None:
                        size_pt = r.font.size.pt
                        break
                break
        except Exception:
            pass
        align = "left"
        try:
            al = shp.text_frame.paragraphs[0].alignment
            if al is not None and "CENTER" in str(al):
                align = "center"
            elif "RIGHT" in str(al):
                align = "right"
        except Exception:
            pass
        _draw_text(d, x, y, w, h, text, size_pt, False, tcolor, align)
    return img


def main():
    prs = Presentation(str(PPTX))
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    files = []
    for i in range(len(prs.slides)):
        img = _render_slide(prs, i)
        p = OUT_DIR / f"slide_{i + 1:02d}.png"
        img.save(p)
        files.append(p)
        print("rendered:", p.name, img.size)
    if TRUNC_WARN:
        print("[WARN] 文本疑似超出文本框（结构预览截断）:")
        for w_ in TRUNC_WARN[:20]:
            print("   ", w_)
    # contact sheet
    thumbs = []
    for p in files:
        im = Image.open(p)
        im.thumbnail((400, 225))
        thumbs.append(im)
    sheet = Image.new("RGB", (400 * 3 + 40, 225 * 5 + 40), (235, 235, 235))
    for i, t in enumerate(thumbs):
        r, c = divmod(i, 3)
        sheet.paste(t, (10 + c * 410, 10 + r * 235))
    sheet.save(OUT_DIR / "contact_sheet.png")
    print("contact sheet:", OUT_DIR / "contact_sheet.png")


if __name__ == "__main__":
    main()
