#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""按方案 A（企业深蓝工程报告）原生重建 7 页可编辑 PPTX。

数据来源：CAW7203 联网报告（2026-08-11 06:00 ~ 08-12 06:00）。
全部文字为原生文本框，图表/表格为 Office 原生元素（可编辑），无图片资产。
"""

from __future__ import annotations

import os
import re
import zipfile
import shutil

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = "/Users/user/Desktop/PY/MC_LogAnalysis/out put/test caw7203/ppt_mockup/A_重制.pptx"

EA_CN = "微软雅黑"
LATIN = "Arial"

# 方案 A 设计系统（参考 UI改善.md 第 7 节 Apple 配色规范）
C_PAGE = "FFFFFF"
C_PANEL = "F5F5F7"
C_PANEL_LINE = "E5E5EA"
C_TITLE = "1D1D1F"
C_SUB = "86868B"
C_TEXT = "3A3A3C"
C_ACCENT = "007AFF"
C_DANGER = "FF3B30"
C_DANGER_LIGHT = "FFB4AB"
C_TAG_BG = "F5F5F7"
C_TAG_FG = "007AFF"
C_TABLE_HEAD = "F9F9FB"
C_TABLE_HEAD_FG = "1D1D1F"
C_TABLE_ALT = "FAFAFA"
C_TABLE_LINE = "E5E5EA"
C_BAR = ["007AFF", "34C759", "FF9F0A"]
C_PIE = ["34C759", "FF9F0A", "FF3B30"]
C_LINE = "007AFF"
C_LINE_HL = "FF3B30"
C_FOOTER = "86868B"
C_GRID = "EDEDF2"
C_AXIS = "86868B"
C_COVER_TOP = "F2F2F7"
C_COVER_BOT = "FFFFFF"
C_COVER_TITLE = "1D1D1F"
C_COVER_SUB = "86868B"

COPYRIGHT = "© 2026 Hon Hai Precision Industry Co., Ltd. All rights reserved."
PX = 13.333 / 600.0  # 基准画布 600px → 13.333in


def P(v):
    return v * PX


def _em(v):
    """转 EMU：已是整型（Length）直接用，浮点按英寸换算，避免二次放大。"""
    return int(v) if isinstance(v, int) else int(round(v * 914400))


def rgb(c):
    return RGBColor.from_string(c)


def set_run(run, text, size, color, bold=False, font_cn=EA_CN):
    run.text = str(text)
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = rgb(color)
    f.name = LATIN
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font_cn)


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: [(text, size, color, bold)] 或 [str]。"""
    tb = slide.shapes.add_textbox(_em(x), _em(y), _em(w), _em(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for line in lines:
        if isinstance(line, str):
            line = (line, 14, C_TEXT, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        set_run(run, line[0], line[1], line[2], line[3])
    return tb


def shape(slide, kind, x, y, w, h, fill=None, line=None, line_w=1.0, radius=None):
    shp = slide.shapes.add_shape(kind, _em(x), _em(y), _em(w), _em(h))
    shp.shadow.inherit = False
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    return shp


def set_fill_opacity(shp, opacity_pct):
    """设置纯色填充透明度（OOXML a:alpha，千分之一百分比）。"""
    srgb = shp.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha = srgb.find(qn("a:alpha"))
    if alpha is None:
        alpha = srgb.makeelement(qn("a:alpha"), {})
        srgb.append(alpha)
    alpha.set("val", str(int(opacity_pct * 1000)))


def line(slide, x0, y0, x1, y1, color, w=1.0):
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, _em(x0), _em(y0), _em(x1), _em(y1)
    )
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(w)
    ln.shadow.inherit = False
    return ln


def chart_font(chart, size=10, color=C_TEXT):
    chart.font.size = Pt(size)
    chart.font.color.rgb = rgb(color)


def fix_chart_axids(chart):
    """把 c:axId / c:crossAx 的负数 val 修正为绝对值（schema 校验要求）。"""
    space = chart._chartSpace
    for tag in ("c:axId", "c:crossAx"):
        for el in space.iter(qn(tag)):
            val = el.get("val")
            if val:
                try:
                    n = int(val)
                except ValueError:
                    continue
                if n < 0:
                    el.set("val", str(-n))


def fix_app_slides(path):
    """docProps/app.xml 的幻灯片数与实际页数同步。"""
    tmp = path + ".tmp"
    with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "docProps/app.xml":
                xml = data.decode("utf-8")
                xml = re.sub(r"<Slides>\d+</Slides>", "<Slides>7</Slides>", xml)
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp, path)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header(slide, idx, total, title, subtitle):
    """内容页页头：tag chip + 标题 + 页码胶囊 + 副标题 + 分隔线。"""
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.53, 0.36, 1.35, 0.42,
          fill=C_TAG_BG, line=None, radius=0.5)
    tb = textbox(slide, 0.53, 0.36, 1.35, 0.42, [("设备能效", 12, C_TAG_FG, True)],
                 anchor=MSO_ANCHOR.MIDDLE)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    textbox(slide, 2.1, 0.26, 8.0, 0.6, [(title, 28, C_TITLE, True)])
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.76, 0.36, 2.04, 0.42,
          fill=C_TAG_BG, line=None, radius=0.5)
    tb = textbox(slide, 10.76, 0.36, 2.04, 0.42,
                 [(f"{idx:02d} / {total:02d}", 12, C_TAG_FG, True)],
                 anchor=MSO_ANCHOR.MIDDLE)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    textbox(slide, 0.53, 0.98, 12.27, 0.42, [(subtitle, 13, C_SUB, False)])
    line(slide, 0.53, 1.52, 12.8, 1.52, C_TITLE, 2)


def footer(slide, idx, total):
    line(slide, 0.53, 7.0, 12.8, 7.0, C_TABLE_LINE, 1)
    textbox(slide, 0.53, 7.1, 7.0, 0.3, [(COPYRIGHT, 9, C_FOOTER, False)])
    tb = textbox(slide, 11.3, 7.1, 1.5, 0.3, [(f"{idx}/{total}", 10, C_ACCENT, True)],
                 anchor=MSO_ANCHOR.MIDDLE)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def panel(slide, x, y, w, h):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
          fill=C_PANEL, line=C_PANEL_LINE, line_w=1.0, radius=0.035)


def make_table(slide, x, y, w, header, rows, col_weights, row_h=0.34, header_h=0.4,
               frame_h=None):
    n = len(rows) + 1
    gf = slide.shapes.add_table(n, len(header), Inches(x), Inches(y),
                                Inches(w),
                                Inches(frame_h or (header_h + row_h * len(rows))))
    tbl = gf.table
    total_w = sum(col_weights)
    for ci, cw in enumerate(col_weights):
        tbl.columns[ci].width = Inches(w * cw / total_w)
    tbl.rows[0].height = Inches(header_h)
    for ri in range(1, n):
        tbl.rows[ri].height = Inches(row_h)
    # python-pptx 设置行高后会重算框高，这里显式写回目标框高（容量/布局按框高判定）
    gf.height = Inches(frame_h or (header_h + row_h * len(rows)))
    for ci, htxt in enumerate(header):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(C_TABLE_HEAD)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        set_run(r, htxt, 11, C_TABLE_HEAD_FG, True)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(C_TABLE_ALT if ri % 2 == 1 else C_PAGE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            r = p.add_run()
            set_run(r, val, 11, C_TEXT, False)
    return tbl


def column_chart_data(cats, vals):
    cd = CategoryChartData()
    cd.categories = list(cats)
    cd.add_series("数值", list(vals))
    return cd


def bar_colors_for(chart, n, hl=None, palette=C_BAR):
    series = chart.plots[0].series[0]
    for i in range(n):
        col = C_ACCENT if (hl is not None and i == hl) else palette[i % len(palette)]
        series.points[i].format.fill.solid()
        series.points[i].format.fill.fore_color.rgb = rgb(col)


def pie_colors_for(chart, n, palette=C_PIE):
    series = chart.plots[0].series[0]
    for i in range(n):
        series.points[i].format.fill.solid()
        series.points[i].format.fill.fore_color.rgb = rgb(palette[i % len(palette)])


def text_len(text, pt):
    """近似文本像素宽（基准 600px 画布）：1pt ≈ 1/1.6px；中文全角、拉丁半角。"""
    unit = pt / 1.6
    return sum(unit if ord(ch) > 127 else unit * 0.55 for ch in str(text))


def wrap_text(text, max_w, pt):
    """按像素宽贪心换行。"""
    lines = []
    cur = ""
    for ch in str(text):
        trial = cur + ch
        if cur and text_len(trial, pt) > max_w:
            lines.append(cur)
            cur = ch
        else:
            cur = trial
    if cur:
        lines.append(cur)
    return lines or [""]


def hbar_list(slide, box, title, rows, color_mode="pareto"):
    """横向条形列表（原生形状+文本框）：标签放得进条形就条内左侧左对齐，
    放不下就在条形右侧换行；数值在远端右对齐；每组垂直居中于条中心。
    rows: [(label, value)] 或 [(label, value, pct)]，pct 存在时数值带「s · %」。
    color_mode='top3'：前 3 条红、其余浅红；color_mode='first'：首条红、其余绿。
    """
    x0, y0, x1, y1 = box
    textbox(slide, P(x0), P(y0 + 2), Inches(4.5), Inches(0.32),
            [(title, 16, C_TITLE, True)])
    py0, py1 = y0 + 20, y1 - 4
    rh = (py1 - py0) / len(rows)
    label_font, value_font = 13, 13
    bar_x0 = x0
    bar_max = (x1 - 150) - bar_x0  # 条形止于数值区之前，避免与数值重叠
    maxv = max(r[1] for r in rows)
    if color_mode == "top3":
        # 按右侧数值判定：第 3 大（去重）为阈值，并列第 3 全部标红
        distinct = sorted({r[1] for r in rows}, reverse=True)
        threshold = distinct[2] if len(distinct) >= 3 else (distinct[-1] if distinct else 0)
    for i, row in enumerate(rows):
        lab, v = row[0], row[1]
        pct = row[2] if len(row) > 2 else None
        cy = py0 + rh * (i + 0.5)
        if color_mode == "top3":
            col = C_DANGER if v >= threshold else C_DANGER_LIGHT
            label_fill = "FFFFFF" if v >= threshold else C_TEXT
        else:
            col = C_DANGER if i == 0 else C_BAR[1]
            label_fill = "FFFFFF"
        bend = bar_x0 + (v / maxv) * bar_max
        bar_w = bend - bar_x0
        shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, P(bar_x0), P(cy - 8),
              P(bend - bar_x0), P(16), fill=col, line=None, radius=0.18)
        value_txt = f"{v:,}s · {pct}%" if pct is not None else f"{v:,}"
        # 数值框：远端右对齐，垂直居中于条形
        tv = textbox(slide, P(x1 - 150), P(cy) - 0.13, Inches(1.5), Inches(0.26),
                     [(value_txt, value_font, C_TEXT, True)], anchor=MSO_ANCHOR.MIDDLE)
        tv.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        single = wrap_text(lab, 99999, label_font)
        fits = bar_w > 48 and len(single) == 1 and text_len(lab, label_font) <= bar_w - 16
        if fits:
            # 标签放得进条形：条内左侧左对齐，垂直居中于条形
            tb = textbox(slide, P(bar_x0 + 8), P(cy) - 0.13, P(bar_w - 16),
                         Inches(0.26), [(lab, label_font, label_fill, True)],
                         anchor=MSO_ANCHOR.MIDDLE)
            tb.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        else:
            # 放不下：标签在条形右侧换行（垂直居中于条形），数值在远端
            value_w = text_len(value_txt, value_font) + 8
            max_w = max(60, (x1 - 6) - (bend + 6) - value_w - 8)
            lines = wrap_text(lab, max_w, label_font)
            shown = lines[:2]
            start_y = cy - (len(shown) * 11) / 2.0
            for k, ln in enumerate(shown):
                textbox(slide, P(bend + 6), P(start_y + k * 11), P(max_w),
                        Inches(0.2), [(ln, label_font, C_TEXT, False)],
                        anchor=MSO_ANCHOR.MIDDLE)


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=C_COVER_TOP)
    bg.fill.gradient()
    stops = bg.fill.gradient_stops
    stops[0].color.rgb = rgb(C_COVER_TOP)
    stops[1].color.rgb = rgb(C_COVER_BOT)
    try:
        bg.fill.gradient_angle = 90
    except Exception:
        pass
    shape(s, MSO_SHAPE.RECTANGLE, 0.53, 3.28, 0.8, 0.13, fill=C_ACCENT)
    textbox(s, 0.53, 3.95, 12.27, 0.9, [("CAW設備能效報告", 40, C_COVER_TITLE, True)])
    textbox(s, 0.53, 5.06, 8.89, 0.55, [("—— UPH、EFF、Alarm", 18, C_COVER_SUB, False)])
    # 磨砂玻璃 KPI 卡（半透明白 + 高光 + 描边）
    kpis = [("18,440", "总产出(个)"), ("768.33", "实际UPH(pcs/h)"), ("62.8%", "综合EFF(%)")]
    gap = 0.28
    total_w = 12.27
    cw = (total_w - gap * (len(kpis) - 1)) / len(kpis)
    for i, (val, lab) in enumerate(kpis):
        x = 0.53 + i * (cw + gap)
        card = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.8, cw, 1.05,
                     fill="FFFFFF", line="D2D2D7", line_w=1.0, radius=0.08)
        set_fill_opacity(card, 35)          # 半透明磨砂
        line(s, x + 0.12, 5.85, x + cw - 0.12, 5.85, "E5E5EA", 0.75)  # 顶部高光
        tb = textbox(s, x, 5.92, cw, 0.55, [(val, 26, C_ACCENT, True)],
                     anchor=MSO_ANCHOR.MIDDLE)
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb = textbox(s, x, 6.48, cw, 0.3, [(lab, 12, C_COVER_SUB, False)],
                     anchor=MSO_ANCHOR.MIDDLE)
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    line(s, 0.53, 6.92, 12.8, 6.92, C_TABLE_LINE, 1)
    textbox(s, 0.53, 7.12, 7.11, 0.31, [(COPYRIGHT, 9, C_COVER_SUB, False)])
    tb = textbox(s, 11.56, 7.12, 1.24, 0.31, [("1/7", 10, C_ACCENT, True)])
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    add_notes(s, "封面：CAW7203 联网分析报告。KPI：总产出 18440、实际UPH 768.33、综合EFF 62.8%。")
    return s


def slide_trend(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 2, 7, "机台状态趋势",
           "每小时 EFF 趋势（低点 2026-08-11 06:00：EFF 0.0%，建议结合停机 Pareto 排查）")
    panel(s, 0.4, 1.72, 12.53, 5.08)
    vals = [0.0, 0.0, 53.17, 89.67, 100.0, 79.03, 72.08, 78.97, 88.81, 100.0, 100.0,
            75.89, 0.0, 0.0, 31.11, 88.39, 94.28, 97.81, 62.33, 94.69, 98.11, 89.28,
            13.64, 0.0]
    cats = [f"{h:02d}:00" for h in range(6, 24)] + ["00:00", "01:00", "02:00",
                                                     "03:00", "04:00", "05:00"]
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.67), Inches(2.3),
                            Inches(12.0), Inches(4.1), column_chart_data(cats, vals))
    ch = gf.chart
    chart_font(ch, 9)
    series = ch.plots[0].series[0]
    series.format.line.color.rgb = rgb(C_LINE)
    series.format.line.width = Pt(2.25)
    series.smooth = False
    ch.plots[0].has_data_labels = False
    ch.value_axis.maximum_scale = 110
    ch.value_axis.minimum_scale = 0
    ch.value_axis.tick_labels.font.size = Pt(8)
    ch.category_axis.tick_labels.font.size = Pt(8)
    ch.category_axis.tick_labels.font.color.rgb = rgb(C_AXIS)
    fix_chart_axids(ch)
    footer(s, 2, 7)
    add_notes(s, "24 小时 EFF 趋势；06:00-08:00 为 0，低点建议结合停机 Pareto 排查。")
    return s


def slide_uph(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 3, 7, "UPH 分析",
           "联网机台产出：总产出 18440 个 · 实际UPH 768.33/h · 1 台机台"
           "（2026/08/11 06:00:00 ~ 2026/08/12 06:00:00）")
    # 上层：柱状图（全宽）
    panel(s, 0.4, 1.72, 12.53, 2.9)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(2.8), Inches(2.05),
                            Inches(7.7), Inches(2.35),
                            column_chart_data(["CAW7203"], [18440]))
    ch = gf.chart
    chart_font(ch, 12)
    ch.plots[0].has_data_labels = True
    ch.plots[0].data_labels.show_value = True
    ch.plots[0].data_labels.number_format = "#,##0"
    ch.plots[0].data_labels.font.size = Pt(14)
    ch.plots[0].data_labels.font.bold = True
    bar_colors_for(ch, 1, hl=None)
    ch.value_axis.maximum_scale = 21000
    fix_chart_axids(ch)
    # 下层：明细表（全宽 7 列）
    panel(s, 0.4, 4.8, 12.53, 2.0)
    make_table(
        s, 0.62, 4.98, 12.1,
        ["机台号", "设备(机型)", "状态", "投入(个)", "产出(个)", "达成率(%)", "UPH(个/小时)"],
        [["CAW7203", "AKC(CAW)", "IDLE", "18587", "18440", "63.6", "1441"]],
        col_weights=[1.0, 1.4, 0.8, 1.0, 1.0, 1.0, 1.2],
        row_h=0.30, header_h=0.42, frame_h=2.0,
    )
    footer(s, 3, 7)
    add_notes(s, "UPH 上下布局：上层各机台产出柱状图，下层 7 列机台明细表。")
    return s


def slide_eff(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 4, 7, "EFF 分析", "EFF = 操作时间(运行+待机) / 计划生产时间")
    panel(s, 0.4, 1.72, 6.5, 5.08)
    gf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.8), Inches(2.2),
                            Inches(5.6), Inches(3.9),
                            column_chart_data(["运行RUN", "待机IDLE", "停机DOWN"],
                                              [46434, 7827, 32139]))
    ch = gf.chart
    chart_font(ch, 12)
    plot = ch.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.show_category_name = False
    plot.data_labels.number_format = "0.0%"
    plot.data_labels.font.size = Pt(12)
    pie_colors_for(ch, 3)
    fix_chart_axids(ch)
    panel(s, 7.3, 1.72, 5.63, 5.08)
    make_table(
        s, 7.48, 1.92, 5.28,
        ["指标", "数值"],
        [
            ["总时间(秒)", "24:00:00"],
            ["计划生产时间(秒)", "24:00:00"],
            ["运行时间RUN(秒)", "12:53:54"],
            ["待机时间IDLE(秒)", "2:10:27"],
            ["停机时间DOWN(秒)", "8:55:39"],
            ["计划停机pDT(秒)", "7:47:29"],
            ["非计划停机uDT(秒)", "1:08:10"],
            ["可用性损失(秒)", "8:55:39"],
            ["操作时间(秒)", "15:04:21"],
            ["EFF(%)", "62.8"],
        ],
        col_weights=[2.1, 1.0],
        row_h=0.38, header_h=0.42,
    )
    footer(s, 4, 7)
    add_notes(s, "EFF 62.8%；时间构成：RUN 46434s / IDLE 7827s / DOWN 32139s。")
    return s


def slide_pareto(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 5, 7, "停机 Pareto",
           "停机主要集中在 常規清潔機台（12 次，占总停机 53.95%），建议优先排查")
    panel(s, 0.4, 1.72, 12.53, 5.08)
    rows = [
        ("常規清潔機台", 17340, "53.95"),
        ("上治具清潔", 10709, "33.32"),
        ("取料2_X1熟料真空檢測異常報警", 2368, "7.37"),
        ("Alam[17,20]^Chassis右取料模組吸真空2", 792, "2.46"),
        ("Alam[16,23]^Chassis左取料模組吸真空連續三次失敗報警", 204, "0.63"),
        ("彈窗.生料倉底部有料感應異常", 104, "0.32"),
        ("Alam[16,17]^Chassis左取料模組吸真空1", 80, "0.25"),
    ]
    hbar_list(s, (30, 110, 570, 298), "停机时长（秒）", rows, color_mode="top3")
    footer(s, 5, 7)
    add_notes(s, "停机 Pareto（横向条形图，全宽）；常規清潔機台 17340s 占 53.95%。")
    return s


def slide_alarm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 6, 7, "报警分析", "按关键词与模块统计报警")
    panel(s, 0.4, 1.72, 12.53, 5.08)
    rows = [
        ("常規清潔機台", 12),
        ("上治具清潔", 10),
        ("取料2_X1熟料真空檢測異常報警", 8),
        ("Alam[17,20]^Chassis右取料模組吸真空2", 8),
        ("取料2.取緩存台不成功報警", 3),
        ("彈窗.生料倉底部有料感應異常", 2),
        ("Alam[17,17]^Chassis右取料模組吸真空1", 2),
        ("Alam[16,20]^Chassis左取料模組吸真空2", 1),
        ("Alam[13,14]^滑台4左上治具吸真空連續抛料報警", 1),
        ("NG盤2_3穴位感應器異常報警", 1),
    ]
    hbar_list(s, (30, 110, 570, 298), "报警次数 Top10", rows, color_mode="top3")
    footer(s, 6, 7)
    add_notes(s, "报警 Top10（横向条形图，全宽），原因名过长时在图表轴自动折行。")
    return s


def slide_status(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, 7, 7, "机台状态汇总", "正常运行占比 53.7%（行业基准 ≥75%）🔴")
    panel(s, 0.4, 1.72, 6.5, 5.08)
    gf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.8), Inches(2.2),
                            Inches(5.6), Inches(3.9),
                            column_chart_data(["RUN", "IDLE", "DOWN"],
                                              [46434, 32139, 7827]))
    ch = gf.chart
    chart_font(ch, 12)
    plot = ch.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.show_category_name = False
    plot.data_labels.number_format = "0.0%"
    plot.data_labels.font.size = Pt(12)
    pie_colors_for(ch, 3)
    fix_chart_axids(ch)
    panel(s, 7.3, 1.72, 5.63, 5.08)
    make_table(
        s, 7.48, 1.92, 5.28,
        ["状态", "次数", "总时长(秒)", "占比(%)"],
        [
            ["RUN", "302", "12:53:54", "53.74"],
            ["DOWN", "58", "8:55:39", "37.2"],
            ["IDLE", "56", "2:10:27", "9.06"],
        ],
        col_weights=[0.9, 0.7, 1.4, 0.9],
        row_h=0.42, header_h=0.42,
    )
    footer(s, 7, 7)
    add_notes(s, "正常运行占比 53.7%，低于行业基准 ≥75%（红）。")
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_trend(prs)
    slide_uph(prs)
    slide_eff(prs)
    slide_pareto(prs)
    slide_alarm(prs)
    slide_status(prs)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    fix_app_slides(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
